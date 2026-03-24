"""
Generate thesis defense presentation slides.
Thesis: Uncertainty Quantification for GNN Surrogates of Agent-Based Transport Models
Author: Mohd Zamin Quadri
Meeting: March 30, 2026 with Dominik Fuchsgruber & Elena Natterer
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──────────────────────────────────────────────────────────────────────
REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR = os.path.join(REPO, "thesis", "latex_tum_official", "figures")
OUT_PATH = os.path.join(REPO, "thesis", "presentation", "thesis_defense.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)


def fig(name):
    """Return full path to a PNG figure."""
    path = os.path.join(FIG_DIR, name)
    if not os.path.exists(path):
        print(f"  WARNING: figure not found: {name}")
        return None
    return path


# ── Colour palette (TUM corporate) ────────────────────────────────────────────
TUM_BLUE = RGBColor(0, 101, 189)  # #0065BD
TUM_DARK = RGBColor(0, 51, 89)  # #003359
TUM_LIGHT = RGBColor(152, 198, 234)  # #98C6EA
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GREY = RGBColor(51, 51, 51)
MED_GREY = RGBColor(102, 102, 102)
LIGHT_GREY = RGBColor(230, 230, 230)
RED_ACCENT = RGBColor(196, 49, 47)
GREEN_ACCENT = RGBColor(56, 142, 60)

# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helper functions ──────────────────────────────────────────────────────────


def add_bg_rect(slide, left, top, width, height, color):
    """Add a filled rectangle as background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_size=18,
    bold=False,
    color=BLACK,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    font_size=18,
    color=DARK_GREY,
    spacing=Pt(6),
    font_name="Calibri",
    bold_prefix=True,
):
    """Add bullet points to a slide. Each bullet can be 'prefix: rest' for bold prefix."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        # Check for bold prefix pattern "**text** rest"
        if bold_prefix and ": " in bullet and not bullet.startswith("("):
            prefix, rest = bullet.split(": ", 1)
            run1 = p.add_run()
            run1.text = prefix + ": "
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = font_name
        else:
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = color
            run.font.name = font_name

    return txBox


def add_slide_number(slide, num, total):
    """Add slide number in bottom-right corner."""
    add_textbox(
        slide,
        Inches(12.0),
        Inches(7.05),
        Inches(1.2),
        Inches(0.35),
        f"{num}/{total}",
        font_size=11,
        color=MED_GREY,
        alignment=PP_ALIGN.RIGHT,
    )


def add_section_header(slide, title, subtitle=""):
    """Create a section divider slide."""
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, TUM_BLUE)
    add_textbox(
        slide,
        Inches(1.5),
        Inches(2.5),
        Inches(10),
        Inches(1.5),
        title,
        font_size=40,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.LEFT,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(1.5),
            Inches(4.0),
            Inches(10),
            Inches(1.0),
            subtitle,
            font_size=22,
            color=TUM_LIGHT,
            alignment=PP_ALIGN.LEFT,
        )


def make_content_slide(title, slide_num, total):
    """Create a standard content slide with title bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # Title bar
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), TUM_BLUE)
    add_textbox(
        slide,
        Inches(0.6),
        Inches(0.15),
        Inches(12),
        Inches(0.8),
        title,
        font_size=28,
        bold=True,
        color=WHITE,
    )
    # Bottom line
    add_bg_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.05), TUM_BLUE)
    add_slide_number(slide, slide_num, total)
    return slide


TOTAL_SLIDES = 22

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, TUM_DARK)
# Accent bar
add_bg_rect(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(0.06), TUM_LIGHT)

# Title
add_textbox(
    slide,
    Inches(1.0),
    Inches(2.2),
    Inches(11.3),
    Inches(1.8),
    "Uncertainty Quantification for Graph Neural\nNetwork Surrogates of Agent-Based\nTransport Models",
    font_size=34,
    bold=True,
    color=WHITE,
    alignment=PP_ALIGN.LEFT,
)

# Author info
add_textbox(
    slide,
    Inches(1.0),
    Inches(4.5),
    Inches(11),
    Inches(0.5),
    "Mohd Zamin Quadri",
    font_size=22,
    bold=True,
    color=TUM_LIGHT,
    alignment=PP_ALIGN.LEFT,
)
add_textbox(
    slide,
    Inches(1.0),
    Inches(5.1),
    Inches(11),
    Inches(0.5),
    "M.Sc. Mathematics in Science and Engineering",
    font_size=16,
    color=TUM_LIGHT,
    alignment=PP_ALIGN.LEFT,
)
add_textbox(
    slide,
    Inches(1.0),
    Inches(5.6),
    Inches(11),
    Inches(0.8),
    "Examiner: Prof. Dr. Stephan Guennemann\n"
    "Advisors: Dominik Fuchsgruber M.Sc., Elena Natterer M.Sc.",
    font_size=14,
    color=WHITE,
    alignment=PP_ALIGN.LEFT,
)
add_textbox(
    slide,
    Inches(1.0),
    Inches(6.5),
    Inches(11),
    Inches(0.4),
    "Data Analytics and Machine Learning (DAML) Group, TUM",
    font_size=14,
    color=MED_GREY,
    alignment=PP_ALIGN.LEFT,
)

# TUM Logo
logo_path = fig("tum_logo.png")
if logo_path:
    slide.shapes.add_picture(logo_path, Inches(10.5), Inches(0.3), height=Inches(1.2))

add_slide_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Outline
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Outline", 2, TOTAL_SLIDES)
sections = [
    "1.  Motivation & Problem Statement",
    "2.  Data & Architecture",
    "3.  Methodology: UQ Methods",
    "4.  Predictive Performance (Trials 1-8)",
    "5.  MC Dropout Uncertainty Quality",
    "6.  Calibration & Temperature Scaling",
    "7.  Conformal Prediction",
    "8.  Selective Prediction",
    "9.  Proper Scoring Rules (CRPS, PIT, Winkler)",
    "10. Ensemble Investigation (Negative Result)",
    "11. Cross-Trial Validation & S-Convergence",
    "12. Key Findings & Contributions",
    "13. Limitations & Future Work",
]
add_bullet_slide(
    slide,
    Inches(1.0),
    Inches(1.4),
    Inches(11),
    Inches(5.8),
    sections,
    font_size=19,
    color=DARK_GREY,
    bold_prefix=False,
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Motivation
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Motivation: Why Uncertainty Matters", 3, TOTAL_SLIDES)
bullets = [
    "Agent-based transport simulations (MATSim): hours per scenario",
    "GNN surrogates: seconds per scenario, but how reliable?",
    "Urban planners need to know where predictions may fail",
    "Without UQ: planners may unknowingly base infrastructure decisions on unreliable predictions",
    "This thesis: systematic UQ evaluation for GNN traffic surrogates",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.5),
    Inches(5.5),
    bullets,
    font_size=19,
    color=DARK_GREY,
    bold_prefix=False,
)

# Figure: with/without UQ
fpath = fig("fig6_with_without_uq.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.6), Inches(1.5), width=Inches(5.2))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Research Questions
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Research Questions", 4, TOTAL_SLIDES)
rqs = [
    "RQ1: How effectively does MC Dropout capture epistemic uncertainty in GNN-based traffic surrogates?",
    "RQ2: How does MC Dropout compare to ensemble-based UQ in terms of uncertainty quality and computational cost?",
    "RQ3: Does combining uncertainty from architecturally identical models provide benefits over single-model UQ?",
    "RQ4: Can distribution-free methods (conformal prediction, post-hoc calibration) transform raw uncertainty into trustworthy prediction intervals?",
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, rq in enumerate(rqs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    # Bold the RQ label
    label, rest = rq.split(": ", 1)
    r1 = p.add_run()
    r1.text = label + ": "
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = TUM_BLUE
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(20)
    r2.font.color.rgb = DARK_GREY
    r2.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Data & Problem
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Data: Paris Road Network", 5, TOTAL_SLIDES)
bullets = [
    "Source: MATSim agent-based simulation (Natterer et al. 2025)",
    "10,000 policy scenarios; this thesis uses 1,000 (10% subset)",
    "Each scenario = one static graph (no temporal dimension)",
    "31,635 nodes (road segments), 5 features per node",
    "Target: change in traffic volume (veh/h) after policy intervention",
    "Split: 800 train / 100 val / 100 test",
    "Test set: 100 graphs x 31,635 nodes = 3,163,500 predictions",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.0),
    Inches(5.5),
    bullets,
    font_size=17,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("fig_network_intro.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.3), Inches(1.4), width=Inches(5.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Model Architecture: PointNetTransfGAT", 6, TOTAL_SLIDES)
bullets = [
    "2 PointNetConv layers (512, 128 channels): encode geometry",
    "2 TransformerConv layers (256, 512 channels, 4 heads): long-range attention",
    "2 GATConv layers (64, 1 channel): final embedding + prediction",
    "Dropout (p=0.2 for T8) applied throughout for MC Dropout UQ",
    "8 trial configurations tested (varying dropout, LR, batch size, data split)",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(5.5),
    Inches(3.0),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("fig8_architecture.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(0.8), Inches(4.2), width=Inches(11.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Section divider - Methodology
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Methodology", "Six UQ approaches evaluated")
add_slide_number(slide, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: UQ Methods Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("UQ Methods: Overview", 8, TOTAL_SLIDES)

methods = [
    "MC Dropout (S=30): multiple stochastic forward passes at inference; mean = prediction, std = uncertainty",
    "Deep Ensembles: train M independent models, use variance across predictions",
    "Combined Uncertainty: merge MC Dropout + ensemble variance",
    "Split Conformal Prediction: distribution-free coverage guarantee via calibration quantile",
    "Selective Prediction: rank by uncertainty, retain most confident subset",
    "Temperature Scaling: post-hoc single-parameter calibration of predictive distribution",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(11.5),
    Inches(5.5),
    methods,
    font_size=17,
    color=DARK_GREY,
    bold_prefix=True,
    spacing=Pt(10),
)

# MC Dropout diagram
fpath = fig("fig13_mc_dropout_inference.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(8.5), Inches(4.5), width=Inches(4.3))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Predictive Performance
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Predictive Performance: Trial Comparison", 9, TOTAL_SLIDES)

# Key numbers
bullets = [
    "T1 (no dropout, Linear head): R2 = 0.7860, MAE = 2.97 veh/h -- but no UQ possible",
    "T8 (best UQ-compatible): R2 = 0.5957, MAE = 3.96, RMSE = 7.12 veh/h",
    "T3-T4 (weighted loss): R2 = 0.22-0.24 -- weighted MSE did not help",
    "T8 selected for all UQ experiments (best among dropout-enabled trials)",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.0),
    Inches(3.0),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("fig12_trial_progression.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.0), Inches(1.3), width=Inches(5.8))

# Bottom note
add_textbox(
    slide,
    Inches(0.8),
    Inches(6.3),
    Inches(11),
    Inches(0.5),
    "Note: Elena's original work (full 10k dataset) achieved R2 = 0.91; this thesis uses 10% subset.",
    font_size=13,
    color=MED_GREY,
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: MC Dropout Uncertainty Quality
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide(
    "MC Dropout Uncertainty Quality (T8, S=30)", 10, TOTAL_SLIDES
)

bullets = [
    "Spearman rho = 0.4820: moderate positive correlation between predicted uncertainty and actual error",
    "AUROC = 0.7585 for top-10% error detection",
    "Mean uncertainty: 1.37 veh/h (std across 30 forward passes)",
    "Per-graph rho: mean = 0.4643, 95% CI [0.4599, 0.4689] (bootstrap, 100 graphs)",
    "Stable across all 100 test scenarios (no graph-level outliers)",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.0),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("fig2_uq_ranking.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.0), Inches(1.3), width=Inches(5.8))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Calibration -- Reliability Diagram
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide(
    "Calibration: Reliability Diagram & Temperature Scaling", 11, TOTAL_SLIDES
)

bullets = [
    "Raw MC Dropout sigma is NOT calibrated: ECE = 0.265",
    "Naive Gaussian 95% interval: only 55.6% coverage (should be 95%)",
    "k_95 = 11.34 (need 11.34 sigma, not 1.96 sigma, for 95%)",
    "Temperature scaling (T = 2.70): ECE drops from 0.269 to 0.048 (82% improvement)",
    "Residual gap: scaled 95% interval achieves only 83.3% coverage",
    "Conclusion: ranking is useful, absolute magnitudes require correction",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.5),
    Inches(4.5),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("t8_temperature_scaling.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.8), Inches(1.3), width=Inches(5.0))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Conformal Prediction
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide(
    "Conformal Prediction: Coverage Guarantees", 12, TOTAL_SLIDES
)

bullets = [
    "Split conformal: 50 calibration + 50 evaluation graphs",
    "90% level: q = 9.92 veh/h, coverage = 90.02%",
    "95% level: q = 14.68 veh/h, coverage = 95.01%",
    "Audit split (20/80): corroborates at q = 14.77, coverage = 95.1%",
    "Distribution-free: no Gaussian assumption needed",
    "Fixed-width limitation: same interval for all nodes",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.5),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("fig14_conformal_workflow.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.0), Inches(1.3), width=Inches(5.8))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Adaptive Conformal
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide(
    "Adaptive Conformal: Node-Specific Intervals", 13, TOTAL_SLIDES
)

bullets = [
    "Standard conformal: fixed-width -- over-covers low-uncertainty nodes (98.6%), under-covers high-uncertainty nodes (62.9%)",
    "Adaptive conformal: normalise residuals by MC Dropout sigma",
    "Result: conditional coverage narrows to [90.0%, 96.2%] across deciles",
    "MC Dropout ranking signal carries genuine calibration information",
    "Practical: wider intervals where model is uncertain, narrower where confident",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.5),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("t8_conformal_conditional.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.0), Inches(1.3), width=Inches(5.8))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Selective Prediction
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide(
    "Selective Prediction: Trading Coverage for Accuracy", 14, TOTAL_SLIDES
)

bullets = [
    "Sort predictions by ascending MC Dropout uncertainty",
    "Retain only the most confident fraction:",
    "  100%: MAE = 3.95 veh/h (baseline)",
    "  90%: MAE = 3.23 veh/h (-18.3%)",
    "  50%: MAE = 2.32 veh/h (-41.2%)",
    "  10%: MAE = 1.16 veh/h (-70.6%)",
    "Monotone improvement confirms uncertainty ranking is meaningful",
    "T7 replication: same pattern (-38.3% at 50%)",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.0),
    Inches(5.0),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("t8_selective_prediction_curve.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.0), Inches(1.5), width=Inches(5.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: Proper Scoring Rules
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Proper Scoring Rules: CRPS, PIT, Winkler", 15, TOTAL_SLIDES)

bullets = [
    "CRPS = 3.38 veh/h; CRPS/MAE = 0.857 (theoretical optimum 0.707 for perfect Gaussian)",
    "21% above optimum: quantifies calibration cost of underdispersed sigma",
    "PIT histogram: spike at 0 (28.4%), U-shaped tails -- systematic underdispersion",
    "After temperature scaling: KS drops 0.245 to 0.104 (-57%)",
    "Winkler scores: conformal beats naive Gaussian (32.3 vs 49.7 at 90% level)",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.5),
    Inches(3.5),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("t8_pit_after_tempscaling.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(1.0), Inches(4.5), width=Inches(5.5))

fpath = fig("t8_interval_width_comparison.png")
if fpath:
    slide.shapes.add_picture(fpath, Inches(7.0), Inches(4.5), width=Inches(5.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: Ensemble Investigation
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide(
    "Ensemble Investigation: A Documented Negative Result", 16, TOTAL_SLIDES
)

bullets = [
    "5 independently trained models: all produced R2 ~ 0.003 when loaded",
    "Root cause: PyTorch Geometric GATConv API version mismatch",
    "Checkpoint stores lin.weight (old PyG); current PyG 2.3.1 expects lin_src.weight / lin_dst.weight",
    "strict=False silently drops old keys -- random initialisation of final 2 layers",
    "Within degraded context: MC Dropout rho = 0.16 vs ensemble rho = 0.10",
    "Standalone MC Dropout (correct PyG version): rho = 0.48 -- unaffected",
    "Conclusion: ensemble experiments are inconclusive (RQ3), not a failure of ensembling itself",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(11.5),
    Inches(5.5),
    bullets,
    font_size=17,
    color=DARK_GREY,
    bold_prefix=False,
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: Cross-Trial & S-Convergence
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide(
    "Robustness: Cross-Trial Validation & S-Convergence", 17, TOTAL_SLIDES
)

bullets = [
    "T7 replication: rho = 0.4460, selective prediction -38.3% at 50%, conformal 95.3% coverage",
    "Same qualitative conclusions hold across both trials",
    "S-convergence: S=30 rho = 0.4584, S=50 rho = 0.4632 (<1% improvement)",
    "Curve flattens sharply after S ~ 25: S=30 is cost-optimal",
    "Inference time: 228 min total (S=30 x 100 graphs) -- manageable overhead",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(6.5),
    Inches(4.5),
    bullets,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

fpath = fig("t8_s_convergence.pdf")
# s_convergence might only be PDF, check for PNG
fpath_png = (
    fig("t8_s_convergence.png")
    if os.path.exists(os.path.join(FIG_DIR, "t8_s_convergence.png"))
    else None
)
if fpath_png:
    slide.shapes.add_picture(fpath_png, Inches(7.5), Inches(1.5), width=Inches(5.3))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18: Section divider - Findings
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Key Findings & Contributions", "Answering RQ1-RQ4")
add_slide_number(slide, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19: RQ Answers
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Answers to Research Questions", 19, TOTAL_SLIDES)

answers = [
    (
        "RQ1",
        "Yes",
        "MC Dropout captures meaningful uncertainty: rho = 0.4820, selective prediction -41.2% MAE at 50% retention, AUROC = 0.76",
    ),
    (
        "RQ2",
        "Yes",
        "MC Dropout outperforms ensemble variance (rho 0.16 vs 0.10 in degraded context), at lower cost. Standalone rho = 0.48",
    ),
    (
        "RQ3",
        "Inconclusive",
        "PyG loading bug invalidated all ensemble results. Identical-architecture ensembles remain untested under correct conditions",
    ),
    (
        "RQ4",
        "Yes",
        "Conformal prediction: 95.01% coverage. Temperature scaling: 82% ECE reduction. Adaptive conformal: [90.0%, 96.2%] conditional coverage",
    ),
]

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.8))
tf = txBox.text_frame
tf.word_wrap = True

for i, (rq, verdict, detail) in enumerate(answers):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    # RQ label
    r1 = p.add_run()
    r1.text = f"{rq}: "
    r1.font.size = Pt(18)
    r1.font.bold = True
    r1.font.color.rgb = TUM_BLUE
    r1.font.name = "Calibri"
    # Verdict
    r2 = p.add_run()
    r2.text = f"{verdict}. "
    r2.font.size = Pt(18)
    r2.font.bold = True
    if verdict == "Inconclusive":
        r2.font.color.rgb = RED_ACCENT
    else:
        r2.font.color.rgb = GREEN_ACCENT
    r2.font.name = "Calibri"
    # Detail
    r3 = p.add_run()
    r3.text = detail
    r3.font.size = Pt(16)
    r3.font.color.rgb = DARK_GREY
    r3.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20: The UQ Hierarchy
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide(
    "Practical Recommendation: Three-Level UQ Hierarchy", 20, TOTAL_SLIDES
)

levels = [
    "Level 1 -- MC Dropout (S=30, dropout=0.2): per-node uncertainty ranking for screening unreliable predictions. Use for selective prediction.",
    "Level 2 -- Temperature Scaling (T=2.70): post-hoc calibration, reduces ECE by 82%. Improves absolute sigma interpretation at no cost.",
    "Level 3 -- Conformal Prediction: formal coverage guarantee (95.01%). Use when contractual or operational coverage is required.",
    "Adaptive Conformal: combines Levels 1+3 for node-specific intervals with both ranking and coverage properties.",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.5),
    Inches(11.5),
    Inches(4.5),
    levels,
    font_size=18,
    color=DARK_GREY,
    bold_prefix=True,
    spacing=Pt(14),
)

# Bottom takeaway
add_bg_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), LIGHT_GREY)
add_textbox(
    slide,
    Inches(0.8),
    Inches(6.1),
    Inches(12.0),
    Inches(0.7),
    "Takeaway: All three levels are achievable with the PointNetTransfGAT architecture at no meaningful cost to predictive accuracy.",
    font_size=17,
    bold=True,
    color=TUM_DARK,
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21: Limitations & Future Work
# ═══════════════════════════════════════════════════════════════════════════════
slide = make_content_slide("Limitations & Future Work", 21, TOTAL_SLIDES)

limitations = [
    "Single network (Paris) -- generalisation to other cities untested",
    "10% data subset (1,000 / 10,000 scenarios) -- full dataset may improve R2",
    "No non-GNN baselines for comparison",
    "Ensemble experiments invalidated by PyG bug",
    "MC Dropout provides ranking, not calibrated probabilities (without post-hoc correction)",
]
add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.4),
    Inches(5.5),
    Inches(3.5),
    limitations,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)

# Future work column
future = [
    "Full 10,000-scenario training",
    "Diverse ensembles (GAT + GraphSAGE + GCN)",
    "Conformal risk control (adaptive coverage)",
    "Post-hoc GEBM (Fuchsgruber et al. 2024)",
    "Multi-city validation",
    "Epistemic-aleatoric decomposition",
]
add_textbox(
    slide,
    Inches(7.0),
    Inches(1.3),
    Inches(5.5),
    Inches(0.5),
    "Future Directions",
    font_size=20,
    bold=True,
    color=TUM_BLUE,
)
add_bullet_slide(
    slide,
    Inches(7.0),
    Inches(1.9),
    Inches(5.5),
    Inches(4.5),
    future,
    font_size=16,
    color=DARK_GREY,
    bold_prefix=False,
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22: Thank You
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, TUM_DARK)
add_bg_rect(slide, Inches(0), Inches(2.2), SLIDE_W, Inches(0.06), TUM_LIGHT)

add_textbox(
    slide,
    Inches(1.0),
    Inches(2.6),
    Inches(11),
    Inches(1.2),
    "Thank You",
    font_size=48,
    bold=True,
    color=WHITE,
    alignment=PP_ALIGN.LEFT,
)

add_textbox(
    slide,
    Inches(1.0),
    Inches(4.0),
    Inches(11),
    Inches(0.6),
    "Questions & Discussion",
    font_size=28,
    color=TUM_LIGHT,
    alignment=PP_ALIGN.LEFT,
)

# Key numbers summary
summary_lines = [
    "MC Dropout rho = 0.4820  |  Selective prediction: -41.2% MAE at 50%  |  Conformal: 95.01% coverage",
    "Temperature scaling: 82% ECE reduction  |  CRPS/MAE = 0.857  |  S=30 is cost-optimal",
]
add_textbox(
    slide,
    Inches(1.0),
    Inches(5.2),
    Inches(11),
    Inches(1.0),
    "\n".join(summary_lines),
    font_size=15,
    color=MED_GREY,
    alignment=PP_ALIGN.LEFT,
)

add_textbox(
    slide,
    Inches(1.0),
    Inches(6.5),
    Inches(11),
    Inches(0.5),
    "Mohd Zamin Quadri  |  M.Sc. Mathematics in Science and Engineering  |  TUM",
    font_size=14,
    color=MED_GREY,
    alignment=PP_ALIGN.LEFT,
)

add_slide_number(slide, 22, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUT_PATH)
print(f"Presentation saved: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
file_size = os.path.getsize(OUT_PATH)
print(f"File size: {file_size / 1024 / 1024:.1f} MB")
