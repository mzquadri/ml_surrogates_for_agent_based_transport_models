#!/usr/bin/env python3
"""
Generate thesis defense presentation.
Thesis: Uncertainty Quantification for GNN Surrogates of Agent-Based Transport Models
Author: Mohd Zamin Quadri, TUM 2025

All HD plots embedded, detailed content, TUM colors.
"""

import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = str(Path(__file__).resolve().parents[1])
PLOTS_DIR = os.path.join(REPO, "docs", "hd_plots")
OUT_PPTX = os.path.join(REPO, "thesis_presentation_final.pptx")

# TUM Colors
TUM_BLUE = RGBColor(0x00, 0x65, 0xBD)
TUM_DARK = RGBColor(0x00, 0x33, 0x59)
TUM_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TUM_LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
TUM_ORANGE = RGBColor(0xE3, 0x72, 0x22)
TUM_GREEN = RGBColor(0xA2, 0xAD, 0x00)
TUM_RED = RGBColor(0xCC, 0x00, 0x33)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide(title_text=None, layout_idx=6):
    """Add a blank slide with optional title."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    return slide


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a TUM blue title bar at the top."""
    # Blue header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TUM_BLUE
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = TUM_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.65), Inches(12), Inches(0.4)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p2.font.italic = True


def add_content_text(
    slide,
    text,
    left=0.5,
    top=1.3,
    width=12,
    height=5.5,
    font_size=16,
    bold=False,
    color=DARK_GRAY,
    alignment=PP_ALIGN.LEFT,
):
    """Add content text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold if not line.startswith("  ") else False
        p.alignment = alignment
        p.space_after = Pt(4)

    return txBox


def add_image(slide, img_name, left=0.3, top=1.3, width=None, height=None):
    """Add an image from hd_plots directory."""
    img_path = os.path.join(PLOTS_DIR, img_name)
    if not os.path.exists(img_path):
        print(f"  WARNING: {img_path} not found!")
        return
    if width and height:
        slide.shapes.add_picture(
            img_path, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    elif width:
        slide.shapes.add_picture(
            img_path, Inches(left), Inches(top), width=Inches(width)
        )
    elif height:
        slide.shapes.add_picture(
            img_path, Inches(left), Inches(top), height=Inches(height)
        )
    else:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top))


def add_footer(
    slide,
    text="M. Z. Quadri | TUM 2025 | Uncertainty Quantification for GNN Surrogates",
):
    """Add footer bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TUM_DARK
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = TUM_WHITE
    p.alignment = PP_ALIGN.LEFT


def add_bullet_text(
    slide, items, left=0.5, top=1.3, width=12, height=5.5, font_size=16
):
    """Add bulleted text."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if item.startswith("##"):
            p.text = item[2:].strip()
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = TUM_BLUE
            p.font.bold = True
            p.space_before = Pt(12)
        elif item.startswith("  -"):
            p.text = item[3:].strip()
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = DARK_GRAY
            p.level = 1
        elif item.startswith("-"):
            p.text = item[1:].strip()
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY
        else:
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY

        p.space_after = Pt(3)

    return txBox


print("Building presentation...")

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = add_slide()

# Large blue background
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
)
shape.fill.solid()
shape.fill.fore_color.rgb = TUM_BLUE
shape.line.fill.background()

# Title
add_content_text(
    slide,
    "Uncertainty Quantification for\nGraph Neural Network Surrogates of\nAgent-Based Transport Models",
    left=1,
    top=1.0,
    width=11.3,
    height=3,
    font_size=36,
    bold=True,
    color=TUM_WHITE,
    alignment=PP_ALIGN.CENTER,
)

# Author info
add_content_text(
    slide,
    "Mohd Zamin Quadri\nM.Sc. Mathematics in Science and Engineering",
    left=1,
    top=3.8,
    width=11.3,
    height=1.2,
    font_size=22,
    bold=False,
    color=RGBColor(0xCC, 0xDD, 0xEE),
    alignment=PP_ALIGN.CENTER,
)

# Supervisors
add_content_text(
    slide,
    "Examiner: Prof. Dr. Stephan Gunnemann\nAdvisors: Dominik Fuchsgruber M.Sc., Elena Natterer M.Sc.\n\nTechnical University of Munich\nDepartment of Informatics | DAML Group | 2025",
    left=1,
    top=5.0,
    width=11.3,
    height=2,
    font_size=16,
    bold=False,
    color=RGBColor(0xAA, 0xCC, 0xDD),
    alignment=PP_ALIGN.CENTER,
)

print("  Slide 1: Title")

# ============================================================
# SLIDE 2: Motivation & Problem
# ============================================================
slide = add_slide()
add_title_bar(slide, "Motivation & Problem Statement")
add_bullet_text(
    slide,
    [
        "## The Challenge",
        "- Agent-based transport simulations (MATSim) are computationally expensive",
        "  - One scenario evaluation takes hours of simulation time",
        "  - Transport planners need to evaluate hundreds of policy scenarios",
        "  - Real-time policy exploration is impossible with traditional simulations",
        "",
        "## Our Approach",
        "- Train a GNN surrogate model to replace expensive simulations",
        "  - Input: graph of Munich road network (31,635 nodes, 5 features each)",
        "  - Output: predicted change in traffic volume (delta_v in veh/h) per road link",
        "  - Inference time: seconds vs hours",
        "",
        "## The Missing Piece: Uncertainty Quantification",
        "- Predictions alone are not enough --- we need confidence estimates",
        "  - Which predictions should planners trust?",
        "  - How to flag unreliable predictions?",
        "  - Can we provide guaranteed coverage intervals?",
    ],
    font_size=15,
)
add_footer(slide)
print("  Slide 2: Motivation")

# ============================================================
# SLIDE 3: Data Overview
# ============================================================
slide = add_slide()
add_title_bar(slide, "Data Overview", "Static, cross-sectional graph data from MATSim")
add_bullet_text(
    slide,
    [
        "## Dataset Characteristics",
        "- 100 test graphs (transport policy scenarios for Munich)",
        "- Each graph: 31,635 nodes (road links)",
        "- Total test nodes: 3,163,500 (100 x 31,635)",
        "- No temporal dimension --- each graph is a static snapshot",
        "",
        "## Input Features (5 per node)",
        "- VOL_BASE_CASE: base traffic volume",
        "- CAPACITY_BASE_CASE: road capacity",
        "- FREESPEED: free-flow speed",
        "- LENGTH: road link length",
        "- CAPACITY_REDUCTION: policy-induced capacity change",
        "",
        "## Target Variable",
        "- delta_v (vehicles/hour): change in traffic volume due to policy",
        "- Range varies significantly across nodes and scenarios",
        "",
        "## Train/Val/Test Split",
        "- 80/10/10 graph-level split",
        "- All UQ evaluation on 100 held-out test graphs",
    ],
    font_size=14,
)
add_footer(slide)
print("  Slide 3: Data Overview")

# ============================================================
# SLIDE 4: Model Architecture
# ============================================================
slide = add_slide()
add_title_bar(slide, "GNN Architecture: Hybrid Multi-Layer Design (Trial 8)")
add_bullet_text(
    slide,
    [
        "## Architecture: PointNetConv + TransformerConv + GATConv",
        "",
        "- Layer 1: PointNetConv (7 -> 256 -> 512) --- Local neighborhood aggregation",
        "- Layer 2: PointNetConv (514 -> 256 -> 128) --- Deeper local features",
        "- Layer 3: TransformerConv (128 -> 256, 4 heads) --- Attention-based message passing",
        "- Layer 4: TransformerConv (256 -> 512, 4 heads) --- Higher-level attention",
        "- Layer 5: GATConv (512 -> 64) --- Graph attention refinement",
        "- Layer 6: GATConv (64 -> 1) --- Final prediction output",
        "",
        "## Key Hyperparameters",
        "- Dropout: p = 0.2 (lower than default; 'lower_dropout' variant)",
        "- Batch size: 8 with gradient accumulation (effective batch = 24)",
        "- Learning rate: 0.0005",
        "",
        "## Deterministic Performance (T8)",
        "- R-squared = 0.5957",
        "- MAE = 3.96 veh/h",
        "- RMSE = 7.12 veh/h",
        "- No overfitting (val/test gap < 0.3%)",
    ],
    font_size=14,
)
add_footer(slide)
print("  Slide 4: Architecture")

# ============================================================
# SLIDE 5: MC Dropout Method
# ============================================================
slide = add_slide()
add_title_bar(slide, "MC Dropout: Core Uncertainty Quantification Method")
add_bullet_text(
    slide,
    [
        "## Method",
        "- Enable dropout at test time (model.train() with frozen BatchNorm)",
        "- Run S = 30 stochastic forward passes per input",
        "- Prediction = mean of 30 outputs",
        "- Uncertainty (sigma) = standard deviation of 30 outputs",
        "",
        "## MC Dropout Performance",
        "- R-squared = 0.5857 (slightly lower than deterministic 0.5957 due to averaging)",
        "- MAE = 3.95 veh/h | RMSE = 7.21 veh/h",
        "- Mean uncertainty sigma = 1.37 veh/h",
        "- Total inference time: ~228 minutes for 100 graphs x 30 samples",
        "",
        "## Uncertainty Quality",
        "- Spearman correlation (sigma vs |error|): rho = 0.4820",
        "  - Moderate positive correlation: higher uncertainty = higher actual error",
        "  - Per-graph rho: mean = 0.4643, 95% CI = [0.4599, 0.4689]",
        "  - Useful for selective prediction and risk flagging",
        "- Bootstrap confidence interval from 10,000 resamples across 100 graphs",
    ],
    font_size=14,
)
add_footer(slide)
print("  Slide 5: MC Dropout")

# ============================================================
# SLIDE 6: S-Convergence Analysis
# ============================================================
slide = add_slide()
add_title_bar(slide, "S-Convergence: Justifying S = 30 MC Samples")
add_image(slide, "02_s_convergence_analysis.png", left=0.3, top=1.3, width=12.5)
add_footer(slide)
print("  Slide 6: S-Convergence")

# ============================================================
# SLIDE 7: Model Performance Dashboard
# ============================================================
slide = add_slide()
add_title_bar(slide, "Model Performance Dashboard: Complete Results Summary")
add_image(
    slide,
    "18_model_performance_dashboard.png",
    left=0.3,
    top=1.2,
    width=12.7,
    height=5.6,
)
add_footer(slide)
print("  Slide 7: Dashboard")

# ============================================================
# SLIDE 8: Selective Prediction
# ============================================================
slide = add_slide()
add_title_bar(slide, "Selective Prediction: Risk-Coverage Trade-off")
add_image(
    slide, "01_selective_prediction_risk_coverage.png", left=0.5, top=1.3, width=7.5
)
add_bullet_text(
    slide,
    [
        "## Key Results",
        "- 100% retained: MAE = 3.95",
        "- 90% retained: MAE = 3.23",
        "  (-18.3% reduction)",
        "- 50% retained: MAE = 2.32",
        "  (-41.2% reduction)",
        "- 25% retained: MAE = 1.79",
        "  (-54.6% reduction)",
        "- 10% retained: MAE = 1.06",
        "",
        "## Practical Impact",
        "- Planners can reject",
        "  uncertain predictions",
        "- Even retaining 90%",
        "  significantly improves",
        "  accuracy",
    ],
    left=8.3,
    top=1.3,
    width=4.5,
    font_size=13,
)
add_footer(slide)
print("  Slide 8: Selective Prediction")

# ============================================================
# SLIDE 9: PIT Before vs After
# ============================================================
slide = add_slide()
add_title_bar(slide, "PIT Calibration: Before vs After Temperature Scaling")
add_image(slide, "05_pit_before_vs_after_comparison.png", left=0.3, top=1.3, width=12.7)
add_footer(slide)
print("  Slide 9: PIT Comparison")

# ============================================================
# SLIDE 10: Reliability Diagram
# ============================================================
slide = add_slide()
add_title_bar(slide, "Reliability Diagram: Coverage Calibration Improvement")
add_image(slide, "06_reliability_diagram.png", left=2.5, top=1.2, width=8.3, height=5.8)
add_footer(slide)
print("  Slide 10: Reliability Diagram")

# ============================================================
# SLIDE 11: Temperature Scaling Summary
# ============================================================
slide = add_slide()
add_title_bar(slide, "Temperature Scaling: Calibration Results")
add_image(slide, "16_temperature_scaling_summary.png", left=0.3, top=1.3, width=8)
add_bullet_text(
    slide,
    [
        "## Temperature Scaling",
        "- Method: sigma_scaled = sigma * T",
        "- Optimal T = 2.7025",
        "  (fitted on 20-graph cal set)",
        "",
        "## Improvements",
        "- ECE: 0.269 -> 0.048 (-82%)",
        "- KS stat: 0.245 -> 0.104 (-57%)",
        "- NLL: 21.6 -> 4.75 (-78%)",
        "- PIT 1st bin: 0.284 -> 0.088",
        "",
        "## Why T = 2.70 is large?",
        "- Raw MC Dropout sigma is",
        "  overconfident (too narrow)",
        "- T > 1 widens the intervals",
        "- After scaling, intervals match",
        "  actual error distribution",
    ],
    left=8.5,
    top=1.3,
    width=4.5,
    font_size=13,
)
add_footer(slide)
print("  Slide 11: Temperature Scaling")

# ============================================================
# SLIDE 12: CRPS by Decile
# ============================================================
slide = add_slide()
add_title_bar(slide, "CRPS Analysis: Scoring Rule by Uncertainty Decile")
add_image(slide, "07_crps_by_decile.png", left=0.3, top=1.3, width=12.7)
add_footer(slide)
print("  Slide 12: CRPS by Decile")

# ============================================================
# SLIDE 13: CRPS/MAE Ratio
# ============================================================
slide = add_slide()
add_title_bar(slide, "CRPS/MAE Ratio: Probabilistic Model Quality Assessment")
add_image(slide, "15_crps_mae_ratio.png", left=0.5, top=1.3, width=7.5)
add_bullet_text(
    slide,
    [
        "## CRPS/MAE Ratio = 0.857",
        "",
        "- CRPS mean = 3.383 veh/h",
        "- MAE = 3.948 veh/h",
        "- Ratio = CRPS/MAE = 0.857",
        "",
        "## Interpretation",
        "- Theoretical optimum: 0.707",
        "  (perfect Gaussian: 1/sqrt(2))",
        "- Our ratio: 0.857",
        "  (~21% above optimum)",
        "- 14.3% below MAE",
        "  (CRPS < MAE = good sign)",
        "",
        "## Meaning",
        "- Probabilistic model provides",
        "  sharper intervals than naive",
        "  point prediction + global spread",
    ],
    left=8.3,
    top=1.3,
    width=4.5,
    font_size=13,
)
add_footer(slide)
print("  Slide 13: CRPS/MAE")

# ============================================================
# SLIDE 14: Winkler Score
# ============================================================
slide = add_slide()
add_title_bar(slide, "Winkler Score: Interval Quality Assessment")
add_image(slide, "08_winkler_score_comparison.png", left=0.3, top=1.3, width=12.7)
add_footer(slide)
print("  Slide 14: Winkler Score")

# ============================================================
# SLIDE 15: Conformal Prediction
# ============================================================
slide = add_slide()
add_title_bar(slide, "Conformal Prediction: Distribution-Free Coverage Guarantees")
add_image(slide, "17_conformal_prediction_intervals.png", left=0.3, top=1.3, width=7.5)
add_bullet_text(
    slide,
    [
        "## Method",
        "- Split data 50/50: calibration/test",
        "- Compute nonconformity scores",
        "- Select quantile for target level",
        "",
        "## Results (absolute intervals)",
        "- 90% target -> 90.0% actual",
        "  (q_hat = 9.92 veh/h)",
        "- 95% target -> 95.0% actual",
        "  (q_hat = 14.68 veh/h)",
        "",
        "## Key Property",
        "- Model-agnostic guarantee",
        "- No distributional assumptions",
        "- Valid even if model is misspecified",
        "- Intervals are wider than",
        "  raw Gaussian (but correct!)",
    ],
    left=8.3,
    top=1.3,
    width=4.5,
    font_size=13,
)
add_footer(slide)
print("  Slide 15: Conformal Prediction")

# ============================================================
# SLIDE 16: Conditional Coverage
# ============================================================
slide = add_slide()
add_title_bar(slide, "Conditional Coverage: Global vs Adaptive Conformal Prediction")
add_image(slide, "09_conformal_conditional_coverage.png", left=0.3, top=1.3, width=12.7)
add_footer(slide)
print("  Slide 16: Conditional Coverage")

# ============================================================
# SLIDE 17: MC vs Conformal Coverage
# ============================================================
slide = add_slide()
add_title_bar(slide, "Raw MC Dropout vs Conformal: Coverage Gap")
add_image(
    slide, "19_mc_dropout_vs_conformal_coverage.png", left=0.5, top=1.3, width=7.5
)
add_bullet_text(
    slide,
    [
        "## The Problem",
        "- Raw Gaussian (MC Dropout):",
        "  90% target -> only 48.6% actual",
        "  95% target -> only 54.8% actual",
        "  (SEVERE under-coverage)",
        "",
        "## Root Cause",
        "- MC Dropout sigma is too narrow",
        "- k95 = 11.34 (needs 11.34 sigma",
        "  for 95% coverage vs 1.96 for",
        "  true Gaussian)",
        "",
        "## The Fix",
        "- Temperature Scaling (T=2.70)",
        "  improves ECE from 0.27 to 0.05",
        "- Conformal Prediction gives",
        "  EXACT coverage guarantees",
        "  without assumptions",
    ],
    left=8.3,
    top=1.3,
    width=4.5,
    font_size=13,
)
add_footer(slide)
print("  Slide 17: MC vs Conformal")

# ============================================================
# SLIDE 18: NLL Comparison
# ============================================================
slide = add_slide()
add_title_bar(slide, "Negative Log-Likelihood: Cross-Model Comparison")
add_image(slide, "10_nll_comparison.png", left=1.5, top=1.3, width=10)
add_footer(slide)
print("  Slide 18: NLL Comparison")

# ============================================================
# SLIDE 19: Per-Graph Distributions
# ============================================================
slide = add_slide()
add_title_bar(slide, "Per-Graph Analysis: Rho, MAE, and Sigma Distributions")

# Add three plots side by side (smaller)
add_image(
    slide, "11_per_graph_rho_distribution.png", left=0.2, top=1.2, width=4.3, height=2.9
)
add_image(
    slide, "12_per_graph_mae_distribution.png", left=4.6, top=1.2, width=4.3, height=2.9
)
add_image(
    slide,
    "13_per_graph_sigma_distribution.png",
    left=9.0,
    top=1.2,
    width=4.1,
    height=2.9,
)

add_bullet_text(
    slide,
    [
        "- Rho is stable across graphs: mean = 0.464, 95% CI = [0.460, 0.469] (narrow spread)",
        "- MAE varies 2.4 - 6.1 veh/h: some policy scenarios are harder to predict than others",
        "- Sigma varies 0.9 - 2.2 veh/h: model correctly assigns more uncertainty to harder scenarios",
    ],
    left=0.3,
    top=4.3,
    width=12.5,
    font_size=14,
)
add_footer(slide)
print("  Slide 19: Per-Graph Distributions")

# ============================================================
# SLIDE 20: Stratified UQ
# ============================================================
slide = add_slide()
add_title_bar(slide, "Stratified UQ: Performance by Traffic Volume Quartile")
add_image(slide, "20_stratified_uq_by_volume.png", left=0.3, top=1.3, width=12.7)
add_footer(slide)
print("  Slide 20: Stratified UQ")

# ============================================================
# SLIDE 21: All Features Stratified
# ============================================================
slide = add_slide()
add_title_bar(slide, "Stratified UQ: Spearman Rho Across All Feature Quartiles")
add_image(slide, "21_stratified_uq_all_features.png", left=0.3, top=1.3, width=12.7)
add_footer(slide)
print("  Slide 21: All Features Stratified")

# ============================================================
# SLIDE 22: Ensemble Bug
# ============================================================
slide = add_slide()
add_title_bar(slide, "Ensemble Experiments: PyG API Mismatch (Documented Bug)")
add_image(slide, "14_ensemble_bug_diagnostic.png", left=0.3, top=1.3, width=8)
add_bullet_text(
    slide,
    [
        "## Root Cause",
        "- PyG GATConv API changed",
        "  between versions",
        "- Old: lin.weight (single)",
        "- New: lin_src.weight +",
        "  lin_dst.weight (split)",
        "- strict=False silently",
        "  drops mismatched weights",
        "",
        "## Impact",
        "- Ensemble R-sq = 0.003",
        "  (effectively random)",
        "- Standalone NPZ results",
        "  are UNAFFECTED (produced",
        "  with correct PyG version)",
        "",
        "## Lesson",
        "- Always verify state_dict",
        "  loading in PyTorch/PyG",
    ],
    left=8.5,
    top=1.3,
    width=4.5,
    font_size=13,
)
add_footer(slide)
print("  Slide 22: Ensemble Bug")

# ============================================================
# SLIDE 23: T7 vs T8
# ============================================================
slide = add_slide()
add_title_bar(slide, "Trial 7 vs Trial 8: Impact of Dropout Rate")
add_image(slide, "23_t7_vs_t8_comparison.png", left=0.3, top=1.3, width=12.7)
add_footer(slide)
print("  Slide 23: T7 vs T8")

# ============================================================
# SLIDE 24: Verification Summary
# ============================================================
slide = add_slide()
add_title_bar(slide, "Numeric Verification: 39/39 Checks PASS")
add_image(
    slide, "22_verification_summary_table.png", left=0.5, top=1.2, width=12, height=5.8
)
add_footer(slide)
print("  Slide 24: Verification Summary")

# ============================================================
# SLIDE 25: Key Contributions
# ============================================================
slide = add_slide()
add_title_bar(slide, "Key Contributions & Findings")
add_bullet_text(
    slide,
    [
        "## 1. UQ Framework for GNN Transport Surrogates",
        "- First systematic UQ study for GNN surrogates of agent-based transport models",
        "- 6 UQ methods + 6 scoring rules applied and evaluated",
        "",
        "## 2. Practical Selective Prediction",
        "- Retaining 50% most certain predictions reduces MAE by 41.2%",
        "- Actionable for transport planners: flag and reject unreliable predictions",
        "",
        "## 3. Calibration Improvement Pipeline",
        "- Temperature scaling (T=2.70) reduces ECE by 82%, NLL by 78%",
        "- Conformal prediction provides distribution-free coverage guarantees",
        "",
        "## 4. Rigorous Evaluation & Reproducibility",
        "- 39/39 numeric checks verified against raw JSON artifacts",
        "- All results traceable: code -> NPZ data -> JSON metrics -> thesis text",
        "- Bootstrap confidence intervals quantify statistical uncertainty",
        "",
        "## 5. Documented Pitfalls",
        "- Ensemble PyG API bug transparently documented (R-sq = 0.003)",
        "- S-convergence analysis justifies S=30 choice (< 1% gain at S=50)",
    ],
    font_size=14,
)
add_footer(slide)
print("  Slide 25: Contributions")

# ============================================================
# SLIDE 26: Limitations & Future Work
# ============================================================
slide = add_slide()
add_title_bar(slide, "Limitations & Future Work")
add_bullet_text(
    slide,
    [
        "## Limitations",
        "- Only MC Dropout explored as primary UQ method (ensemble had PyG bug)",
        "- Spearman rho = 0.48 is moderate (not strong) uncertainty-error correlation",
        "- No spatial/graph-structure-aware uncertainty analysis (e.g., road centrality)",
        "- Temperature scaling is global (single T for all nodes); node-level calibration could improve",
        "- R-squared = 0.59 leaves room for base model improvement",
        "",
        "## Future Directions",
        "- Graph-aware UQ methods (e.g., GEBM from Fuchsgruber et al., NeurIPS 2024)",
        "- Spatially adaptive temperature scaling (different T per subgraph region)",
        "- Temporal extensions: time-varying transport scenarios",
        "- Ensemble methods with corrected PyG API (fix weight remapping)",
        "- Conditional coverage analysis with more sophisticated stratification",
        "- Integration with real-time transport planning decision systems",
        "- Multi-task prediction: volume + speed + travel time simultaneously",
    ],
    font_size=14,
)
add_footer(slide)
print("  Slide 26: Limitations")

# ============================================================
# SLIDE 27: Summary
# ============================================================
slide = add_slide()

# Blue background
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
)
shape.fill.solid()
shape.fill.fore_color.rgb = TUM_BLUE
shape.line.fill.background()

add_content_text(
    slide,
    "Summary",
    left=1,
    top=0.8,
    width=11,
    height=1,
    font_size=36,
    bold=True,
    color=TUM_WHITE,
    alignment=PP_ALIGN.CENTER,
)

add_content_text(
    slide,
    "MC Dropout (S=30) provides moderate but useful uncertainty estimates\n"
    "(Spearman rho = 0.48) for GNN surrogates of transport simulations.\n\n"
    "Selective prediction reduces MAE by 41.2% when retaining the 50%\n"
    "most certain predictions (3.95 -> 2.32 veh/h).\n\n"
    "Temperature scaling (T=2.70) dramatically improves calibration:\n"
    "ECE reduced by 82%, NLL by 78%, KS statistic by 57%.\n\n"
    "Conformal prediction provides distribution-free coverage guarantees:\n"
    "90% target -> 90.0% actual, 95% target -> 95.0% actual.\n\n"
    "All 39 numeric claims verified against raw data artifacts.",
    left=1.5,
    top=2.0,
    width=10.3,
    height=4.5,
    font_size=18,
    bold=False,
    color=TUM_WHITE,
    alignment=PP_ALIGN.CENTER,
)

add_content_text(
    slide,
    "Thank you! Questions?",
    left=1,
    top=6.2,
    width=11,
    height=1,
    font_size=28,
    bold=True,
    color=RGBColor(0xCC, 0xDD, 0xEE),
    alignment=PP_ALIGN.CENTER,
)

print("  Slide 27: Summary")

# ============================================================
# SAVE
# ============================================================
prs.save(OUT_PPTX)
print(f"\nPresentation saved: {OUT_PPTX}")
print(f"Total slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(OUT_PPTX) / (1024 * 1024):.1f} MB")
