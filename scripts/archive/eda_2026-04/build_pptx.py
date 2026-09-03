"""
build_pptx.py
Generates presentation_uq.pptx — a faithful Python-pptx version of presentation_uq.tex
TUM Blue colour scheme, 16:9, ~25 slides
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Constants ────────────────────────────────────────────────────────────────
W = Inches(13.33)  # 16:9 width
H = Inches(7.5)  # 16:9 height

TUM_BLUE = RGBColor(0, 101, 189)
TUM_DK = RGBColor(0, 51, 89)
TUM_LIGHT = RGBColor(152, 198, 234)
TUM_GRAY = RGBColor(88, 88, 90)
TUM_GREEN = RGBColor(0, 124, 48)
TUM_RED = RGBColor(196, 7, 27)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
LIGHT_BG = RGBColor(235, 245, 255)  # very light blue for content boxes

FIGS = "slides_figs"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helper functions ──────────────────────────────────────────────────────────


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0.75)):
    shp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l,
        t,
        w,
        h,
    )
    shp.line.width = line_width
    if fill is None:
        shp.fill.background()
        shp.line.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line_color:
            shp.line.color.rgb = line_color
        else:
            shp.line.fill.background()
    return shp


def textbox(
    slide,
    text,
    l,
    t,
    w,
    h,
    font_size=Pt(13),
    bold=False,
    color=BLACK,
    align=PP_ALIGN.LEFT,
    wrap=True,
    italic=False,
):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def title_bar(slide, title_text, subtitle=None):
    """Blue header bar with white title text."""
    bar = rect(slide, 0, 0, W, Inches(1.05), fill=TUM_BLUE)
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    bar.left = 0
    bar.top = 0
    if subtitle:
        textbox(
            slide,
            subtitle,
            Inches(0.2),
            Inches(0.78),
            W - Inches(0.4),
            Inches(0.4),
            font_size=Pt(11),
            color=TUM_LIGHT,
            italic=True,
        )
    return bar


def footer_bar(slide):
    """Dark blue footer."""
    fb = rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill=TUM_DK)
    textbox(
        slide,
        "Mohd Zamin Quadri  |  TUM CIT  |  Master's Thesis Defence 2026",
        Inches(0.2),
        H - Inches(0.33),
        W - Inches(2.5),
        Inches(0.32),
        font_size=Pt(9),
        color=TUM_LIGHT,
    )


def block(
    slide,
    header,
    body_lines,
    l,
    t,
    w,
    h,
    header_color=TUM_BLUE,
    body_bg=LIGHT_BG,
    header_font=Pt(12),
    body_font=Pt(11),
):
    """A coloured block box with header + body text."""
    hh = Inches(0.36)
    # header
    hdr = rect(slide, l, t, w, hh, fill=header_color)
    tf = hdr.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = header
    r.font.size = header_font
    r.font.bold = True
    r.font.color.rgb = WHITE
    # body
    body = rect(slide, l, t + hh, w, h - hh, fill=body_bg, line_color=header_color)
    tf2 = body.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = line
        r2.font.size = body_font
        r2.font.color.rgb = BLACK
    return body


def alert_block(slide, header, body_lines, l, t, w, h, body_font=Pt(11)):
    return block(
        slide,
        header,
        body_lines,
        l,
        t,
        w,
        h,
        header_color=TUM_RED,
        body_bg=RGBColor(255, 235, 235),
        body_font=body_font,
    )


def example_block(slide, header, body_lines, l, t, w, h, body_font=Pt(11)):
    return block(
        slide,
        header,
        body_lines,
        l,
        t,
        w,
        h,
        header_color=TUM_GREEN,
        body_bg=RGBColor(225, 245, 230),
        body_font=body_font,
    )


def add_image(slide, fname, l, t, w, h=None):
    path = os.path.join(FIGS, fname)
    if not os.path.exists(path):
        print(f"  WARNING: {path} not found — skipping image")
        return
    if h is None:
        pic = slide.shapes.add_picture(path, l, t, width=w)
    else:
        pic = slide.shapes.add_picture(path, l, t, width=w, height=h)
    return pic


def table_shape(
    slide,
    data,
    col_widths,
    l,
    t,
    row_height=Inches(0.34),
    header_bg=TUM_BLUE,
    alt_bg=LIGHT_BG,
    font_size=Pt(10.5),
):
    """data: list of rows; first row = header."""
    from pptx.util import Pt

    rows = len(data)
    cols = len(data[0])
    total_w = sum(col_widths)
    tbl = slide.shapes.add_table(rows, cols, l, t, total_w, row_height * rows).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ri, row in enumerate(data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = font_size
            if ri == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                run.font.color.rgb = BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
    return tbl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
# dark blue background band top third
top = rect(sl, 0, 0, W, Inches(3.5), fill=TUM_DK)
# white main area
rect(sl, 0, Inches(3.5), W, H - Inches(3.5), fill=WHITE)

textbox(
    sl,
    "Uncertainty Quantification for\nGNN Surrogates of Agent-Based Transport Models",
    Inches(0.6),
    Inches(0.6),
    W - Inches(1.2),
    Inches(2.2),
    font_size=Pt(30),
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)

textbox(
    sl,
    "Mohd Zamin Quadri",
    Inches(0.6),
    Inches(3.65),
    W - Inches(1.2),
    Inches(0.55),
    font_size=Pt(20),
    bold=True,
    color=TUM_BLUE,
    align=PP_ALIGN.CENTER,
)

textbox(
    sl,
    "TUM School of Computation, Information and Technology",
    Inches(0.6),
    Inches(4.2),
    W - Inches(1.2),
    Inches(0.4),
    font_size=Pt(14),
    color=TUM_GRAY,
    align=PP_ALIGN.CENTER,
)

textbox(
    sl,
    "Advisors: Dominik Fuchsgruber, M.Sc.  &  Elena Natterer, M.Sc.\n"
    "Supervisor: Prof. Dr. Stephan Günnemann",
    Inches(0.6),
    Inches(4.65),
    W - Inches(1.2),
    Inches(0.7),
    font_size=Pt(12),
    color=TUM_GRAY,
    align=PP_ALIGN.CENTER,
)

textbox(
    sl,
    "Master's Thesis Defence — 2026",
    Inches(0.6),
    Inches(5.55),
    W - Inches(1.2),
    Inches(0.4),
    font_size=Pt(13),
    italic=True,
    color=TUM_GRAY,
    align=PP_ALIGN.CENTER,
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Outline
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Presentation Outline")

LX = Inches(0.3)
RX = Inches(6.9)
TY = Inches(1.15)
BW = Inches(6.3)
BH = Inches(5.7)

block(
    sl,
    "UQ Methods Covered",
    [
        "P1  Ranking Uncertainty",
        "     • MC Dropout",
        "     • Selective Prediction",
        "     • Error Detection (AUROC)",
        "",
        "P2  Calibration",
        "     • Why Gaussian intervals fail",
        "     • Temperature Scaling",
        "     • PIT + NLL",
        "",
        "P3  Coverage Guarantees",
        "     • Global Conformal Prediction",
        "     • Adaptive Conformal Prediction",
    ],
    LX,
    TY,
    BW,
    BH,
    body_font=Pt(12),
)

block(
    sl,
    "Further Experiments",
    [
        "P4  Ensemble Experiments",
        "     • Exp A: Same architecture, 5 seeds",
        "     • Exp B: Multi-architecture ensemble",
        "",
        "P5  Negative Results",
        "     • Trial 7: Cross-validation replication",
        "     • Trial 9: Heteroscedastic NLL (failure)",
        "",
        "P6  Summary & Limitations",
    ],
    RX,
    TY,
    BW,
    Inches(4.4),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Total: 13 UQ experiments",
    [
        "All evaluated on Trial 8",
        "Best GATConv model: R² = 0.596",
        "Test set: 100 graphs · 3.16M road-segment nodes",
    ],
    RX,
    TY + Inches(4.5),
    BW,
    Inches(1.15),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Why UQ?
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Motivation: Why Do We Need Uncertainty Quantification?")

LX = Inches(0.3)
TY = Inches(1.15)
BW = Inches(6.3)

block(
    sl,
    "The Problem with Point Predictions",
    [
        "Model predicts:  Δflow = −3.2 veh/hr  on road segment i",
        "A traffic planner acts on this number.",
        "But HOW RELIABLE is this prediction?",
    ],
    LX,
    TY,
    BW,
    Inches(1.6),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Two road segments, same predicted value:",
    [
        "Segment A: model has seen many similar cases  ✓",
        "Segment B: unusual geometry, out-of-distribution  ✗",
        "A point prediction cannot distinguish them.",
    ],
    LX,
    TY + Inches(1.72),
    BW,
    Inches(1.65),
    body_font=Pt(12),
)

example_block(
    sl,
    "What UQ provides:",
    [
        "Output              Value      Use",
        "─────────────────────────────────────",
        "Point pred. ŷ      −3.2       Forecast",
        "Uncertainty σ       0.4       Trust score",
        "Interval        [−13.2, 6.8]  Hard bound",
        "",
        "Goal: Make the model say not just WHAT it",
        "predicts, but HOW CONFIDENT it is.",
    ],
    Inches(6.8),
    TY,
    Inches(6.3),
    Inches(3.5),
    body_font=Pt(11.5),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — UQ Framework 3 Levels
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "UQ Framework: Three Levels of Uncertainty")

# table
tdata = [
    ["Level", "Question answered", "Method"],
    ["1. Ranking", "Which predictions are least reliable?", "MC Dropout σ"],
    ["1. Ranking", "Can we improve quality by abstaining?", "Selective Prediction"],
    ["1. Ranking", "Can σ detect large errors?", "Error Detection (AUROC)"],
    ["2. Calibration", "Are confidence scores accurate?", "Temperature Scaling"],
    ["2. Calibration", "Are probability distributions correct?", "PIT, NLL"],
    ["3. Guarantee", "Guaranteed coverage interval?", "Global Conformal"],
    ["3. Guarantee", "Conditional coverage interval?", "Adaptive Conformal"],
]
table_shape(
    sl,
    tdata,
    [Inches(1.8), Inches(4.8), Inches(2.9)],
    Inches(0.3),
    Inches(1.2),
    row_height=Inches(0.52),
    font_size=Pt(11.5),
)

block(
    sl,
    "Key Insight",
    [
        "Each level adds value:",
        "  Ranking   → operational  (decide when to trust)",
        "  Calibration → diagnostic  (is σ meaningful?)",
        "  Guarantee  → formal       (provable bound)",
    ],
    Inches(9.85),
    Inches(1.2),
    Inches(3.25),
    Inches(2.55),
    body_font=Pt(11),
)

alert_block(
    sl,
    "Model",
    [
        "Trial 8: PointNet + Transformer + GATConv",
        "R² = 0.596,  MC S = 30",
        "Test: 100 graphs,  3.16M nodes",
    ],
    Inches(9.85),
    Inches(3.88),
    Inches(3.25),
    Inches(1.5),
    body_font=Pt(11),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MC Dropout Method
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "MC Dropout: Method", subtitle="Part 1: Ranking Uncertainty")

LX = Inches(0.3)
TY = Inches(1.15)
BW = Inches(6.3)

block(
    sl,
    "Standard Dropout (Training)",
    [
        "Randomly zero out neurons during training",
        "→ acts as regularisation",
    ],
    LX,
    TY,
    BW,
    Inches(1.15),
    body_font=Pt(12),
)

alert_block(
    sl,
    "MC Dropout (Inference) — Gal & Ghahramani 2016",
    [
        "KEEP dropout active at test time.",
        "Run S = 30 forward passes.",
        "",
        "  ŷ = (1/S) Σ f_θs(x)",
        "  σ = std deviation across S passes",
        "",
        "σ = epistemic uncertainty proxy",
    ],
    LX,
    TY + Inches(1.27),
    BW,
    Inches(2.7),
    body_font=Pt(12),
)

example_block(
    sl,
    "Interpretation",
    [
        "High σ → model disagrees with itself → LESS reliable",
        "Low  σ → model is consistent        → MORE reliable",
    ],
    LX,
    TY + Inches(4.1),
    BW,
    Inches(1.25),
    body_font=Pt(12),
)

# right: figure
add_image(sl, "uncertainty_hist.png", Inches(6.9), Inches(1.15), Inches(6.15))
textbox(
    sl,
    "Distribution of σ across 3.16M test nodes.\nHeavy right tail = small fraction of hard cases.",
    Inches(6.9),
    Inches(5.55),
    Inches(6.15),
    Inches(0.65),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    align=PP_ALIGN.CENTER,
    italic=True,
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MC Dropout Results
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "MC Dropout: Results — Is σ Informative?",
    subtitle="Part 1: Ranking Uncertainty",
)

# left: figure
add_image(sl, "binned_error_vs_uncertainty.png", Inches(0.3), Inches(1.15), Inches(5.9))
textbox(
    sl,
    "MAE increases monotonically with uncertainty bins.\nThis is the key diagnostic for ranking validity.",
    Inches(0.3),
    Inches(5.6),
    Inches(5.9),
    Inches(0.65),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    italic=True,
)

# right
RX = Inches(6.5)
TY = Inches(1.15)
block(
    sl,
    "Point Accuracy: Deterministic vs MC",
    [
        "Mode                R²      MAE    RMSE",
        "────────────────────────────────────────",
        "Deterministic      0.5957  3.96   7.12",
        "MC Dropout (S=30)  0.5857  3.95   7.21",
        "",
        "Accuracy is comparable.",
        "MC Dropout adds σ at marginal cost.",
    ],
    RX,
    TY,
    Inches(6.55),
    Inches(2.6),
    body_font=Pt(11.5),
)

alert_block(
    sl,
    "Ranking Result",
    [
        "Spearman ρ = 0.482  between σ and |error|",
        "",
        "Moderate-to-strong correlation.",
        "σ reliably identifies harder predictions.",
        "",
        "Note: σ is a ranking signal, NOT a calibrated probability.",
    ],
    RX,
    TY + Inches(2.72),
    Inches(6.55),
    Inches(2.65),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Selective Prediction
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "Selective Prediction: Abstaining on Uncertain Cases",
    subtitle="Part 1: Ranking Uncertainty",
)

add_image(sl, "risk_coverage_curve.png", Inches(0.3), Inches(1.15), Inches(5.9))
textbox(
    sl,
    "MAE drops as we reject the most uncertain predictions.\nX-axis = fraction of predictions retained.",
    Inches(0.3),
    Inches(5.6),
    Inches(5.9),
    Inches(0.65),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    italic=True,
)

RX = Inches(6.5)
TY = Inches(1.15)
block(
    sl,
    "Method",
    [
        "Sort all 3.16M predictions by σ (descending).",
        "Reject top-k% most uncertain.",
        "Measure MAE on remaining predictions.",
    ],
    RX,
    TY,
    Inches(6.55),
    Inches(1.55),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Results",
    [
        "Retained   MAE (veh/hr)   Reduction",
        "────────────────────────────────────",
        "100% (baseline)   3.94      —",
        " 90%              3.22    −18.5%",
        " 50%              2.31    −41.6%",
        " 25%              1.77    −55.0%",
    ],
    RX,
    TY + Inches(1.67),
    Inches(6.55),
    Inches(2.55),
    body_font=Pt(12),
)

example_block(
    sl,
    "Practical Value",
    [
        "A planner can choose a confidence threshold.",
        "At 90% retention: nearly 1-in-5 improvement in accuracy.",
        "σ is operationally useful, not just diagnostic.",
    ],
    RX,
    TY + Inches(4.34),
    Inches(6.55),
    Inches(1.45),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Error Detection
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "Error Detection: Can σ Flag Large Errors?",
    subtitle="Part 1: Ranking Uncertainty",
)

add_image(sl, "hexbin_uncertainty_error.png", Inches(0.3), Inches(1.15), Inches(5.9))
textbox(
    sl,
    "Joint density of σ vs. absolute error.\nHigh-error nodes (right) have elevated σ (top).",
    Inches(0.3),
    Inches(5.6),
    Inches(5.9),
    Inches(0.65),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    italic=True,
)

RX = Inches(6.5)
TY = Inches(1.15)
block(
    sl,
    "Binary Detection Task",
    [
        "Positive class: predictions in top-10% of errors",
        "Score: MC Dropout σ",
        "Question: Does high σ predict high error?",
    ],
    RX,
    TY,
    Inches(6.55),
    Inches(1.65),
    body_font=Pt(12),
)

alert_block(
    sl,
    "AUROC Results",
    [
        "Threshold      AUROC    AUPRC",
        "────────────────────────────────",
        "Top-10% errors  0.759    0.315",
        "Top-20% errors  0.740    0.455",
        "Random baseline 0.500    0.100",
    ],
    RX,
    TY + Inches(1.77),
    Inches(6.55),
    Inches(2.35),
    body_font=Pt(12),
)

example_block(
    sl,
    "Interpretation",
    [
        "AUROC = 0.759:",
        "σ ranks a bad prediction above a good one 75.9% of the time.",
        "52% above random baseline — strong discrimination.",
    ],
    RX,
    TY + Inches(4.24),
    Inches(6.55),
    Inches(1.45),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Why Gaussian fails
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Why Naïve Gaussian Intervals Fail", subtitle="Part 2: Calibration")

add_image(sl, "coverage_curve_k_sigma.png", Inches(0.3), Inches(1.15), Inches(5.9))
textbox(
    sl,
    "Empirical coverage vs. multiplier k.\nDashed = Gaussian ideal. Actual coverage far below.",
    Inches(0.3),
    Inches(5.6),
    Inches(5.9),
    Inches(0.65),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    italic=True,
)

RX = Inches(6.5)
TY = Inches(1.15)
block(
    sl,
    "Gaussian Assumption",
    [
        "Build interval: ŷ ± k·σ",
        "Assume σ is a calibrated standard deviation.",
        "k = 1.65  →  90% coverage  (under Gaussian).",
    ],
    RX,
    TY,
    Inches(6.55),
    Inches(1.55),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Reality: Severe Under-coverage",
    [
        "Nominal   k       Actual coverage",
        "──────────────────────────────────",
        " 50%     0.67       23.3%",
        " 80%     1.28       40.1%",
        " 90%     1.65       48.6%   ← 90% target!",
        " 95%     1.96       54.8%",
    ],
    RX,
    TY + Inches(1.67),
    Inches(6.55),
    Inches(2.55),
    body_font=Pt(12),
)

example_block(
    sl,
    "Root Cause: Heavy-Tailed Residuals",
    [
        "Empirical k₉₅ = 11.34   vs.   Gaussian z₉₅ = 1.96",
        "Ratio = 5.79×  heavier tails than Gaussian.",
        "σ is a ranking signal, NOT a calibrated std dev.",
    ],
    RX,
    TY + Inches(4.34),
    Inches(6.55),
    Inches(1.45),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Temperature Scaling
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "Temperature Scaling: Calibrating Confidence Scores",
    subtitle="Part 2: Calibration",
)

LX = Inches(0.3)
TY = Inches(1.15)
BW = Inches(6.5)

block(
    sl,
    "Method (Guo et al., 2017)",
    [
        "A single learnable scalar T scales the uncertainties:",
        "",
        "   σ_cal = T · σ_raw,   T > 0",
        "",
        "T is fitted on a held-out calibration set (minimise ECE or NLL).",
        "",
        "  T > 1 : model was overconfident  → inflate σ",
        "  T < 1 : model was underconfident → shrink σ",
        "  T = 1 : already perfectly calibrated",
    ],
    LX,
    TY,
    BW,
    Inches(3.0),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Trial 8 Result:  T = 2.70",
    [
        "Model was SEVERELY overconfident.",
        "Scaling inflated σ by factor 2.70.",
    ],
    LX,
    TY + Inches(3.12),
    BW,
    Inches(1.25),
    body_font=Pt(12),
)

RX = Inches(7.1)
example_block(
    sl,
    "Calibration Improvement",
    [
        "Metric     Before   After    Change",
        "────────────────────────────────────",
        "ECE         0.269   0.048    −82%",
        "NLL        21.65    4.75     −78%",
        "KS (PIT)    0.245   0.104    −57%",
    ],
    RX,
    TY,
    Inches(6.0),
    Inches(2.35),
    body_font=Pt(12),
)

block(
    sl,
    "Residual Limitation",
    [
        "Even after scaling, slight overconfidence persists",
        "at high prediction magnitudes (large |Δflow|).",
        "",
        "→ Motivates Conformal Prediction which gives",
        "  hard guarantees without distributional assumptions.",
    ],
    RX,
    TY + Inches(2.47),
    Inches(6.0),
    Inches(2.1),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PIT + NLL
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Calibration Diagnostics: PIT and NLL", subtitle="Part 2: Calibration")

LX = Inches(0.3)
TY = Inches(1.15)
BW = Inches(6.5)
block(
    sl,
    "PIT — Probability Integral Transform",
    [
        "If predictions ŷ ~ N(μ̂, σ̂²) are correct,",
        "the PIT values  u_i = Φ((y_i − μ̂_i) / σ̂_i)",
        "should be  Uniform[0,1].",
        "",
        "Kolmogorov-Smirnov statistic measures deviation:",
        "   KS = sup_u | F_empirical(u) − u |",
        "",
        "Raw MC:              KS = 0.245",
        "After scaling (T=2.70): KS = 0.104",
        "Improvement:             −57%",
    ],
    LX,
    TY,
    BW,
    Inches(3.5),
    body_font=Pt(12),
)

RX = Inches(7.1)
block(
    sl,
    "NLL — Negative Log-Likelihood",
    [
        "Measures how well the predicted distribution",
        "covers the true targets:",
        "",
        "NLL = mean_i [ (y_i−μ̂_i)² / (2σ̂_i²) + ½ log σ̂_i² ]",
        "",
        "Raw MC:       NLL = 21.65",
        "After scaling: NLL =  4.75",
        "Improvement:       −78%",
    ],
    RX,
    TY,
    Inches(6.0),
    Inches(2.95),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Summary of Calibration Evidence",
    [
        "ECE:    −82% after T = 2.70 scaling",
        "PIT/KS: −57%",
        "NLL:    −78%",
        "All three metrics agree: SUBSTANTIAL improvement.",
    ],
    RX,
    TY + Inches(3.07),
    Inches(6.0),
    Inches(2.15),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conformal Theory
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "Conformal Prediction: Distribution-Free Coverage Guarantees",
    subtitle="Part 3: Coverage Guarantees",
)

LX = Inches(0.3)
TY = Inches(1.15)
block(
    sl,
    "Core Idea (Vovk et al., 2005)",
    [
        "No assumptions about residual distribution.",
        "Only requires: calibration & test data are EXCHANGEABLE.",
        "",
        "Procedure:",
        " 1. Split test nodes into calibration (50%) and evaluation (50%)",
        " 2. Compute residuals on calibration: r_i = |y_i − ŷ_i|",
        " 3. Compute empirical quantile at level 1−α:",
        "      q_α = Quantile_{1−α}({ r_i }_{i ∈ D_cal})",
        " 4. Prediction interval for new point:   ŷ ± q_α",
    ],
    LX,
    TY,
    Inches(6.5),
    Inches(4.2),
    body_font=Pt(12),
)

RX = Inches(7.1)
alert_block(
    sl,
    "Formal Guarantee",
    [
        "P( y ∈ [ŷ − q_α, ŷ + q_α] ) ≥ 1 − α",
        "",
        "This holds REGARDLESS of the error distribution.",
        "No Gaussian assumption.",
        "No calibration quality assumption.",
    ],
    RX,
    TY,
    Inches(6.0),
    Inches(2.3),
    body_font=Pt(12),
)

block(
    sl,
    "Important Caveat",
    [
        "This is MARGINAL coverage — an average over all nodes.",
        "Per-node conditional coverage is not guaranteed without",
        "additional methods (e.g., conformalized quantile regression).",
    ],
    RX,
    TY + Inches(2.42),
    Inches(6.0),
    Inches(1.85),
    body_font=Pt(12),
)

example_block(
    sl,
    "Why better than Gaussian?",
    [
        "Gaussian: assumes σ is calibrated. It is NOT.",
        "Conformal: uses empirical residuals. NO assumption.",
    ],
    RX,
    TY + Inches(4.39),
    Inches(6.0),
    Inches(1.3),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Global Conformal Results
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl, "Global Conformal Prediction: Results", subtitle="Part 3: Coverage Guarantees"
)

add_image(sl, "coverage_comparison.png", Inches(0.3), Inches(1.15), Inches(5.9))
textbox(
    sl,
    "Conformal (blue) achieves nominal coverage.\nGaussian kσ (red) severely under-covers.",
    Inches(0.3),
    Inches(5.6),
    Inches(5.9),
    Inches(0.65),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    italic=True,
)

RX = Inches(6.5)
TY = Inches(1.15)
alert_block(
    sl,
    "Coverage at Nominal Levels",
    [
        "Target   q (half-width)   Coverage   Width",
        "────────────────────────────────────────────",
        " 50%        1.87           50.35%     3.75",
        " 80%        5.97           80.34%    11.94",
        " 90%        9.99           90.17%    19.99  ← used",
        " 95%       14.77           95.09%    29.54",
        "",
        "All widths in veh/hr.",
    ],
    RX,
    TY,
    Inches(6.55),
    Inches(2.85),
    body_font=Pt(12),
)

block(
    sl,
    "Key Comparison  (at 90% target)",
    [
        "Gaussian kσ:     actual coverage = 48.6%",
        "Global Conformal: actual coverage = 90.17%",
        "Difference: +41.6 percentage points",
    ],
    RX,
    TY + Inches(2.97),
    Inches(6.55),
    Inches(1.65),
    body_font=Pt(12),
)

example_block(
    sl,
    "Interpretation",
    [
        "For 1,000 road segments, conformal guarantees",
        "≈900 have true Δflow inside the interval.",
    ],
    RX,
    TY + Inches(4.74),
    Inches(6.55),
    Inches(1.05),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Adaptive Conformal
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "Adaptive Conformal Prediction: Uncertainty-Scaled Intervals",
    subtitle="Part 3: Coverage Guarantees",
)

LX = Inches(0.3)
TY = Inches(1.15)
block(
    sl,
    "Limitation of Global Conformal",
    [
        "Fixed width q_α for ALL predictions.",
        "Confident predictions get the same wide interval",
        "as uncertain ones.",
    ],
    LX,
    TY,
    Inches(6.5),
    Inches(1.6),
    body_font=Pt(12),
)

block(
    sl,
    "Adaptive Method",
    [
        "Scale interval by MC Dropout σ:",
        "",
        "   r_i_scaled = |y_i − ŷ_i| / (σ_i + ε)",
        "",
        "Compute quantile q_α_adapt on scaled residuals.",
        "Final interval: ŷ ± q_α_adapt · σ",
        "",
        "  Narrow where σ is small",
        "  Wide   where σ is large",
    ],
    LX,
    TY + Inches(1.72),
    Inches(6.5),
    Inches(3.2),
    body_font=Pt(12),
)

RX = Inches(7.1)
add_image(sl, "interval_width_comparison.png", RX, TY, Inches(6.0))
textbox(
    sl,
    "Adaptive intervals adapt width to local uncertainty.",
    RX,
    TY + Inches(3.4),
    Inches(6.0),
    Inches(0.4),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    italic=True,
)

alert_block(
    sl,
    "Conditional Coverage Improvement",
    [
        "σ bin      Global    Adaptive",
        "────────────────────────────────",
        "Low  σ      62.9%     90.0%",
        "High σ      98.6%     96.2%",
        "Overall     90.2%     90.1%",
        "",
        "Adaptive narrows the conditional coverage spread",
        "from [62.9%, 98.6%]  to  [90.0%, 96.2%].",
    ],
    RX,
    TY + Inches(3.9),
    Inches(6.0),
    Inches(2.6),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Ensemble Exp A
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "Ensemble Experiment A: Same Architecture, 5 Random Seeds",
    subtitle="Part 4: Ensemble Experiments",
)

add_image(sl, "exp_a_mc_vs_ensemble.png", Inches(0.3), Inches(1.15), Inches(5.9))
textbox(
    sl,
    "MC Dropout σ vs. ensemble variance as uncertainty proxies.",
    Inches(0.3),
    Inches(5.6),
    Inches(5.9),
    Inches(0.45),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    italic=True,
)

RX = Inches(6.5)
TY = Inches(1.15)
block(
    sl,
    "Setup",
    [
        "5 instances of Trial 8, each trained with a different random seed.",
        "Ensemble variance = variance across 5 model predictions.",
        "Compare to MC Dropout σ from single model.",
    ],
    RX,
    TY,
    Inches(6.55),
    Inches(1.65),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Results: Spearman ρ for uncertainty ranking",
    [
        "Method              Spearman ρ",
        "────────────────────────────────",
        "MC Dropout σ           0.491",
        "Ensemble Variance      0.307",
    ],
    RX,
    TY + Inches(1.77),
    Inches(6.55),
    Inches(1.95),
    body_font=Pt(12),
)

block(
    sl,
    "Why ensemble variance is weaker here",
    [
        "Same architecture + same training data → CORRELATED errors.",
        "Models fail TOGETHER on the same hard nodes.",
        "Ensemble variance does not capture cases where all models",
        "are simultaneously wrong.",
        "",
        "MC Dropout with single model outperforms same-arch ensemble.",
    ],
    RX,
    TY + Inches(3.84),
    Inches(6.55),
    Inches(2.0),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Ensemble Exp B
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "Ensemble Experiment B: Multi-Architecture Ensemble",
    subtitle="Part 4: Ensemble Experiments",
)

add_image(sl, "exp_b_model_comparison.png", Inches(0.3), Inches(1.15), Inches(5.9))
textbox(
    sl,
    "Individual trial R² vs. weighted ensemble R².",
    Inches(0.3),
    Inches(5.6),
    Inches(5.9),
    Inches(0.45),
    font_size=Pt(10.5),
    color=TUM_GRAY,
    italic=True,
)

RX = Inches(6.5)
TY = Inches(1.15)
block(
    sl,
    "Setup",
    [
        "Weighted ensemble of 5 architecturally distinct trials:",
        "T2 + T5 + T6 + T7 + T8",
        "Weights proportional to individual R².",
    ],
    RX,
    TY,
    Inches(6.55),
    Inches(1.55),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Result: Ensemble Underperforms Single Best Model",
    [
        "Method                  R²",
        "───────────────────────────",
        "Trial 8 alone (best)   0.596",
        "Multi-arch ensemble    0.566",
        "",
        "Ensemble uncertainty Spearman ρ = 0.433",
        "vs. MC Dropout ρ = 0.491",
    ],
    RX,
    TY + Inches(1.67),
    Inches(6.55),
    Inches(2.55),
    body_font=Pt(12),
)

block(
    sl,
    "Lesson: Why ensembles can fail",
    [
        "• Weaker models DILUTE the best model's signal",
        "• All trained on same 1,000 scenarios → correlated predictions",
        "• Ensemble diversity requires DATA diversity,",
        "  not just architecture diversity",
    ],
    RX,
    TY + Inches(4.34),
    Inches(6.55),
    Inches(1.55),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Trial 7 Cross-validation
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl, "Trial 7: Cross-Validation Replication", subtitle="Part 5: Negative Results"
)

LX = Inches(0.3)
TY = Inches(1.15)
block(
    sl,
    "Purpose",
    [
        "Trial 8 uses an 80/10/10 split with 100 test graphs.",
        "Trial 7 uses the same split but different hyperparameters.",
        "Goal: verify UQ patterns are not artefacts of Trial 8's specific test set.",
    ],
    LX,
    TY,
    Inches(6.5),
    Inches(1.65),
    body_font=Pt(12),
)

block(
    sl,
    "Trial 7 Setup",
    [
        "Architecture: GATConv (same as T8)",
        "Dropout: 0.3 (higher than T8's 0.2)",
        "Test graphs: 100,  LR: 0.0006",
        "R² = 0.547",
    ],
    LX,
    TY + Inches(1.77),
    Inches(6.5),
    Inches(1.95),
    body_font=Pt(12),
)

RX = Inches(7.1)
alert_block(
    sl,
    "Key UQ Metrics — T7 vs T8",
    [
        "Metric            T8       T7",
        "─────────────────────────────────",
        "Spearman ρ       0.482    0.446",
        "k₉₅             11.34   16.15",
        "Conformal q₉₀    9.99    higher",
    ],
    RX,
    TY,
    Inches(6.0),
    Inches(2.35),
    body_font=Pt(12),
)

block(
    sl,
    "Interpretation",
    [
        "• Pattern REPLICATED: MC Dropout still informative (ρ = 0.446)",
        "• k₉₅ = 16.15 vs T8's 11.34:",
        "  higher dropout = more uncertainty spread",
        "• Gaussian intervals even more miscalibrated",
        "• Confirms findings are model-independent",
    ],
    RX,
    TY + Inches(2.47),
    Inches(6.0),
    Inches(2.3),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Trial 9 Heteroscedastic Failure
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(
    sl,
    "Trial 9: Heteroscedastic NLL — A Negative Result",
    subtitle="Part 5: Negative Results",
)

LX = Inches(0.3)
TY = Inches(1.15)
block(
    sl,
    "Motivation",
    [
        "Instead of post-hoc UQ (MC Dropout), train the model to output",
        "both mean AND variance directly.",
        "",
        "Final layer: GATConv(64 → 2) outputs [μ̂, log σ̂²]",
        "",
        "Loss (Kendall & Gal, 2017):",
        "   L = mean_i [ (y_i−μ̂_i)²/(2σ̂_i²) + ½ log σ̂_i² ]",
    ],
    LX,
    TY,
    Inches(6.5),
    Inches(2.95),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Result: R² = 0.02 — Near-Total Failure",
    [
        "The model learned to predict, but made completely unreliable predictions.",
    ],
    LX,
    TY + Inches(3.07),
    Inches(6.5),
    Inches(0.85),
    body_font=Pt(12),
)

RX = Inches(7.1)
alert_block(
    sl,
    "Root Cause: Variance Inflation (Seitzer et al., 2022)",
    [
        "The NLL loss has a shortcut:",
        "",
        "If σ̂² → ∞, the first term → 0",
        "regardless of μ̂.",
        "",
        "The model learned to INFLATE σ̂² to minimise",
        "loss WITHOUT improving μ̂.",
    ],
    RX,
    TY,
    Inches(6.0),
    Inches(2.95),
    body_font=Pt(12),
)

example_block(
    sl,
    "Lesson & Fix",
    [
        "• Standard heteroscedastic NLL is INSUFFICIENT for this task",
        "• Fix requires: variance head constraints,",
        "  evidential deep learning (Amini et al., 2020),",
        "  or decoupled mean/variance training",
        "• Post-hoc UQ (MC Dropout + Conformal) is more robust",
    ],
    RX,
    TY + Inches(3.07),
    Inches(6.0),
    Inches(2.65),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — All 13 Methods Summary Table
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "All UQ Experiments: Complete Summary")

tdata = [
    ["#", "Method", "Key Result", "Verdict"],
    ["1", "MC Dropout (S=30)", "Spearman ρ = 0.482", "✓ Useful ranking"],
    ["2", "Gaussian kσ intervals", "90% nominal → 48.6% actual", "✗ Miscalibrated"],
    ["3", "Temperature Scaling (T=2.70)", "ECE −82%, NLL −78%", "✓ Calibrates σ"],
    ["4", "PIT Diagnostic", "KS 0.245 → 0.104 (−57%)", "✓ Confirms T.S."],
    ["5", "NLL Diagnostic", "21.65 → 4.75", "✓ Confirms T.S."],
    ["6", "Global Conformal", "90.02% @ 90% target", "✓ Guaranteed"],
    ["7", "Adaptive Conformal", "Conditional [90%, 96.2%]", "✓ Better cond."],
    ["8", "Selective Prediction", "−41.6% MAE @ 50% retain", "✓ Actionable"],
    ["9", "Error Detection (AUROC)", "AUROC = 0.759", "✓ Detects errors"],
    ["10", "Ensemble Exp A (5 seeds)", "ρ = 0.307 < MC's 0.491", "✗ Correlated"],
    ["11", "Ensemble Exp B (multi-arch)", "R² = 0.566 < T8's 0.596", "✗ Diluted"],
    ["12", "Trial 7 Replication", "ρ = 0.446, pattern holds", "✓ Cross-validated"],
    ["13", "Trial 9 Heteroscedastic", "R² = 0.02", "✗ Var. inflation"],
]
table_shape(
    sl,
    tdata,
    [Inches(0.55), Inches(3.6), Inches(4.3), Inches(2.8)],
    Inches(0.3),
    Inches(1.15),
    row_height=Inches(0.435),
    font_size=Pt(11),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — 3-Level Hierarchy Recap
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Three-Level UQ Hierarchy: Final Picture")

TY = Inches(1.2)
BW = Inches(4.0)

block(
    sl,
    "Level 1: Ranking",
    [
        '"Which predictions to trust less?"',
        "",
        "MC Dropout:     ρ = 0.482",
        "AUROC:          0.759",
        "Selective: −41.6% MAE at 50%",
        "",
        "Use: Flag unreliable roads before",
        "deploying model output",
    ],
    Inches(0.3),
    TY,
    BW,
    Inches(4.5),
    body_font=Pt(12),
)

block(
    sl,
    "Level 2: Calibration",
    [
        '"Are confidence scores accurate?"',
        "",
        "Temp. Scaling T = 2.70",
        "ECE: −82%",
        "NLL: −78%,  PIT KS: −57%",
        "",
        "Use: Meaningful probability statements",
        "on individual predictions",
    ],
    Inches(4.65),
    TY,
    BW,
    Inches(4.5),
    body_font=Pt(12),
)

block(
    sl,
    "Level 3: Guarantee",
    [
        '"Hard provable bounds?"',
        "",
        "Global conformal:",
        "  90.02% at 90%",
        "  95.09% at 95%",
        "Adaptive: conditional [90%, 96.2%]",
        "",
        "Use: Worst-case planning for",
        "high-stakes decisions",
    ],
    Inches(9.0),
    TY,
    BW,
    Inches(4.5),
    body_font=Pt(12),
)

# arrow bar
arrow = rect(
    sl, Inches(0.3), H - Inches(1.2), W - Inches(0.6), Inches(0.35), fill=TUM_BLUE
)
textbox(
    sl,
    "← Increasing formal rigour →",
    Inches(0.3),
    H - Inches(1.2),
    W - Inches(0.6),
    Inches(0.35),
    font_size=Pt(12),
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Limitations
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Limitations")

LX = Inches(0.3)
TY = Inches(1.15)
BW = Inches(6.4)

alert_block(
    sl,
    "UQ-Specific Limitations",
    [
        "• Marginal vs. conditional coverage:",
        "  Conformal guarantees hold on AVERAGE over all nodes",
        "  — not per individual node",
        "",
        "• Calibration uses test split:",
        "  Test pool was split 50/50 for conformal calibration.",
        "  Standard practice but worth noting.",
        "",
        "• MC Dropout σ is epistemic only:",
        "  Does not capture aleatoric (data) uncertainty",
        "",
        "• Heavy computational cost:",
        "  MC Dropout (S=30): ~228 min vs. deterministic ~3.4 min (67× overhead)",
    ],
    LX,
    TY,
    BW,
    Inches(5.7),
    body_font=Pt(11.5),
)

alert_block(
    sl,
    "Model & Data Limitations",
    [
        "• Data scarcity:",
        "  1,000 of 10,000 available MATSim scenarios (10% subset)",
        "  Best GATConv model: R²=0.596 vs. Elena's R²=0.78 at 10,000 scenarios",
        "",
        "• Architecture gap:",
        "  Trial 1 (Linear head) achieves R²=0.786 — bottleneck may be in GATConv",
        "",
        "• Single city:",
        "  All results on Paris network. Generalisation to other cities untested.",
        "",
        "• Heteroscedastic approach failed:",
        "  Trial 9 R²=0.02. Native uncertainty prediction requires future work.",
    ],
    Inches(6.95),
    TY,
    BW,
    Inches(5.7),
    body_font=Pt(11.5),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Conclusion")

LX = Inches(0.3)
TY = Inches(1.15)
block(
    sl,
    "What Was Achieved",
    [
        "• Applied 13 UQ experiments on a GNN surrogate for Paris traffic",
        "",
        "• Showed that MC Dropout provides a useful uncertainty signal",
        "  (ρ = 0.482, AUROC = 0.759)",
        "",
        "• Demonstrated that Conformal Prediction gives hard coverage guarantees",
        "  (90.02% at 90% target) where Gaussian methods fail (48.6%)",
        "",
        "• Selective prediction delivers actionable accuracy gains",
        "  (−41.6% MAE at 50% retention)",
        "",
        "• Identified NEGATIVE results: ensemble failure due to correlation;",
        "  heteroscedastic NLL failure due to variance inflation",
    ],
    LX,
    TY,
    Inches(6.5),
    Inches(5.7),
    body_font=Pt(12),
)

RX = Inches(7.1)
example_block(
    sl,
    "Key Numbers",
    [
        "Metric                            Value",
        "────────────────────────────────────────",
        "MC Dropout Spearman ρ             0.482",
        "Conformal coverage @ 90%         90.02%",
        "Adaptive conditional range      90–96.2%",
        "Selective MAE (50% retain)       −41.6%",
        "Error detection AUROC             0.759",
        "ECE improvement (Temp. Scaling)   −82%",
        "Temperature scaling T             2.70",
    ],
    RX,
    TY,
    Inches(6.0),
    Inches(3.8),
    body_font=Pt(12),
)

alert_block(
    sl,
    "Core Message",
    [
        "A model that knows what it doesn't know",
        "is more useful than a more accurate model",
        "that doesn't.",
    ],
    RX,
    TY + Inches(3.92),
    Inches(6.0),
    Inches(1.75),
    body_font=Pt(13),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Appendix A: k95 math
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Appendix A: Why k₉₅ = 11.34 Instead of 1.96?")

LX = Inches(0.3)
TY = Inches(1.15)
block(
    sl,
    "Empirical k Definition",
    [
        "Let k_p be the value such that exactly p% of test residuals satisfy:",
        "   |y_i − ŷ_i| ≤ k_p · σ_i",
        "",
        "Under Gaussian assumptions: k₉₅ = z₀.₉₇₅ = 1.96",
        "Under our heavy-tailed residuals: k₉₅ = 11.34",
        "",
        "Ratio = 11.34 / 1.96 = 5.79 ×",
    ],
    LX,
    TY,
    Inches(6.5),
    Inches(2.85),
    body_font=Pt(12),
)

block(
    sl,
    "What This Means",
    [
        "The residuals are much larger relative to σ than Gaussian predicts.",
        "σ UNDERESTIMATES the spread of true errors by a factor of ~6.",
        "This is why ŷ ± 1.96σ only covers 54.8% rather than 95%.",
    ],
    LX,
    TY + Inches(2.97),
    Inches(6.5),
    Inches(1.75),
    body_font=Pt(12),
)

RX = Inches(7.1)
example_block(
    sl,
    "Why Heavy Tails?",
    [
        "• σ captures EPISTEMIC uncertainty (model disagreement)",
        "• True residuals also contain ALEATORIC uncertainty",
        "  (data noise, unseen disruptions)",
        "• MC Dropout σ has no mechanism to estimate",
        "  aleatoric uncertainty",
        "• Therefore σ systematically underestimates",
        "  total uncertainty",
    ],
    RX,
    TY,
    Inches(6.0),
    Inches(3.1),
    body_font=Pt(12),
)

block(
    sl,
    "T7 Comparison",
    [
        "Trial 7 (dropout=0.3) has k₉₅ = 16.15",
        "Higher dropout → more stochastic → larger σ...",
        "but residuals grow even faster.",
        "",
        "Conformal bypasses this entirely",
        "by using empirical residuals.",
    ],
    RX,
    TY + Inches(3.22),
    Inches(6.0),
    Inches(2.3),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Appendix B: Conformal proof sketch
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Appendix B: Why Conformal Coverage is Guaranteed")

LX = Inches(0.3)
TY = Inches(1.15)
block(
    sl,
    "Exchangeability Assumption",
    [
        "Let (X₁,Y₁), …, (Xₙ,Yₙ), (Xₙ₊₁,Yₙ₊₁) be EXCHANGEABLE",
        "— any permutation of indices has the same joint distribution.",
        "",
        "This holds when calibration and test data are i.i.d. draws",
        "from the same distribution.",
    ],
    LX,
    TY,
    Inches(6.5),
    Inches(2.35),
    body_font=Pt(12),
)

block(
    sl,
    "Proof Sketch (Angelopoulos & Bates, 2023)",
    [
        "The calibration residuals r₁, …, rₙ and the test residual rₙ₊₁",
        "are exchangeable.",
        "",
        "Therefore rₙ₊₁ is equally likely to be in any rank",
        "among {r₁, …, rₙ, rₙ₊₁}.",
        "",
        "P(rₙ₊₁ ≤ q_α) ≥ 1−α follows directly from",
        "order statistics of exchangeable variables.",
    ],
    LX,
    TY + Inches(2.47),
    Inches(6.5),
    Inches(3.1),
    body_font=Pt(12),
)

RX = Inches(7.1)
alert_block(
    sl,
    "Finite-Sample Guarantee",
    [
        "With n calibration points:",
        "",
        "  P(Y ∈ C(X)) ≥ 1 − α − 1/(n+1)",
        "",
        "For n = 1,581,750 (50% of test nodes):",
        "The correction term 1/(n+1) ≈ 6×10⁻⁷ is negligible.",
    ],
    RX,
    TY,
    Inches(6.0),
    Inches(2.75),
    body_font=Pt(12),
)

example_block(
    sl,
    "No Distributional Assumptions",
    [
        "• Works for any model",
        "• Works for any residual distribution",
        "• Works whether or not the model is well-calibrated",
        "• Only fails if exchangeability is violated (covariate shift)",
    ],
    RX,
    TY + Inches(2.87),
    Inches(6.0),
    Inches(2.3),
    body_font=Pt(12),
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Appendix C: All trial R² table
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
footer_bar(sl)
title_bar(sl, "Appendix C: All Trial Results Overview")

tdata = [
    ["T#", "Architecture", "Key Change", "Split", "R²", "MAE", "RMSE"],
    ["1", "Linear", "Baseline, bs=32", "80/15/5", "0.786", "2.97", "5.40"],
    ["2", "GATConv", "First GATConv, bs=16", "80/15/5", "0.512", "4.33", "8.15"],
    ["3", "GATConv", "Weighted loss, no dropout", "80/15/5", "0.225", "5.99", "10.27"],
    ["4", "GATConv", "Weighted loss + dropout", "80/15/5", "0.243", "6.08", "10.15"],
    ["5", "GATConv", "Paper config, bs=8", "80/15/5", "0.555", "4.24", "7.78"],
    ["6", "GATConv", "Lower LR (0.0003)", "80/15/5", "0.522", "4.32", "8.06"],
    ["7", "GATConv", "80/10/10, LR=0.0006", "80/10/10", "0.547", "4.06", "7.53"],
    ["8", "GATConv ★", "Lower dropout (0.2)", "80/10/10", "0.596", "3.96", "7.12"],
    ["9", "Heteroscedastic", "NLL loss, output [μ, σ²]", "80/10/10", "0.02", "—", "—"],
]
table_shape(
    sl,
    tdata,
    [
        Inches(0.5),
        Inches(2.0),
        Inches(3.5),
        Inches(1.4),
        Inches(0.85),
        Inches(0.85),
        Inches(0.95),
    ],
    Inches(0.3),
    Inches(1.2),
    row_height=Inches(0.5),
    font_size=Pt(11),
)

alert_block(
    sl,
    "Note on Trial 1",
    [
        "T1 uses LINEAR final layer — architecturally distinct from T2–T8 (GATConv).",
        "T1 R²=0.786 is not directly comparable.",
        "Best GATConv model = T8, R²=0.596.",
    ],
    Inches(0.3),
    Inches(5.85),
    Inches(6.3),
    Inches(1.45),
    body_font=Pt(11.5),
)

block(
    sl,
    "Elena's Reference",
    [
        "Natterer et al. achieved R²=0.78 with 10,000 scenarios (10× more data).",
        "Our T8 uses only 1,000 scenarios (10% subset).",
        "Data gap explains most of the R² difference.",
    ],
    Inches(6.95),
    Inches(5.85),
    Inches(6.1),
    Inches(1.45),
    body_font=Pt(11.5),
)

# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
out = "presentation_uq.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
