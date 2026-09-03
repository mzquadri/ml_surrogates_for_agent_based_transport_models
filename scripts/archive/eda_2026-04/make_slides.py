"""
Thesis Presentation: ML Surrogates for Agent-Based Transport Models
Creates a clean, professional PowerPoint with charts and no overlapping text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Figure directory ──────────────────────────────────────────────────────────
FIG_DIR = r"C:\Users\zamin\Downloads\Nazim\Thesis\Nazim\ml_surrogates_thesis_final\document\figures"
OUT_PATH = r"C:\Users\zamin\Downloads\Nazim\Thesis_Presentation.pptx"

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Brand colours ─────────────────────────────────────────────────────────────
C_DARK_BLUE = RGBColor(0x1A, 0x3A, 0x5C)  # headings / title bg
C_MID_BLUE = RGBColor(0x26, 0x6B, 0xB2)  # accent
C_LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)  # subtle bg
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY = RGBColor(0x60, 0x60, 0x60)
C_GREEN = RGBColor(0x2E, 0x86, 0x48)
C_ORANGE = RGBColor(0xE0, 0x70, 0x20)
C_RED = RGBColor(0xC0, 0x28, 0x28)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank

# ═══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════════════════════


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Pt

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l),
        Inches(t),
        Inches(w),
        Inches(h),
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    l,
    t,
    w,
    h,
    text,
    font_size=Pt(14),
    bold=False,
    italic=False,
    color=C_BLACK,
    align=PP_ALIGN.LEFT,
    word_wrap=True,
):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_textbox_lines(
    slide,
    l,
    t,
    w,
    h,
    lines,
    font_size=Pt(13),
    color=C_BLACK,
    bold=False,
    line_spacing=1.15,
    align=PP_ALIGN.LEFT,
):
    """Add a textbox with multiple lines (list of strings or (text, bold, color) tuples)."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align

        # line can be str or (str, bold_bool, color)
        if isinstance(line, str):
            txt, b, clr = line, bold, color
        else:
            txt = line[0]
            b = line[1] if len(line) > 1 else bold
            clr = line[2] if len(line) > 2 else color

        run = p.add_run()
        run.text = txt
        run.font.size = font_size
        run.font.bold = b
        run.font.color.rgb = clr

        # line spacing
        from pptx.oxml.ns import qn

        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(line_spacing * 100000)))

    return txb


def header_bar(slide, title_text, subtitle_text="", slide_num=None):
    """Dark blue header bar at the top."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=C_DARK_BLUE)
    add_textbox(
        slide,
        0.3,
        0.08,
        11.5,
        0.6,
        title_text,
        font_size=Pt(26),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.LEFT,
    )
    if subtitle_text:
        add_textbox(
            slide,
            0.3,
            0.65,
            10,
            0.4,
            subtitle_text,
            font_size=Pt(14),
            color=C_LIGHT_BLUE,
            align=PP_ALIGN.LEFT,
        )
    if slide_num:
        add_textbox(
            slide,
            12.5,
            0.1,
            0.7,
            0.35,
            str(slide_num),
            font_size=Pt(13),
            color=RGBColor(0xAA, 0xCC, 0xEE),
            align=PP_ALIGN.RIGHT,
        )


def add_figure(slide, fname, l, t, w, h):
    """Safely insert a PNG figure."""
    path = os.path.join(FIG_DIR, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        # placeholder box
        add_rect(
            slide, l, t, w, h, fill_rgb=RGBColor(0xEE, 0xEE, 0xEE), line_rgb=C_GRAY
        )
        add_textbox(
            slide,
            l + 0.05,
            t + h / 2 - 0.2,
            w - 0.1,
            0.4,
            f"[Figure: {fname}]",
            font_size=Pt(10),
            color=C_GRAY,
            align=PP_ALIGN.CENTER,
        )


def caption_box(slide, l, t, w, text, color=C_GRAY, size=Pt(11)):
    add_textbox(
        slide,
        l,
        t,
        w,
        0.4,
        text,
        font_size=size,
        italic=True,
        color=color,
        align=PP_ALIGN.CENTER,
    )


def bullet_section(
    slide,
    l,
    t,
    w,
    h,
    title,
    bullets,
    title_size=Pt(15),
    bullet_size=Pt(13),
    title_color=C_DARK_BLUE,
    bullet_color=C_BLACK,
    bg_color=None,
    border_color=C_MID_BLUE,
):
    """Draws a bordered section box with a title and bullet list."""
    if bg_color:
        add_rect(
            slide,
            l,
            t,
            w,
            h,
            fill_rgb=bg_color,
            line_rgb=border_color,
            line_width=Pt(1.5),
        )
    add_textbox(
        slide,
        l + 0.1,
        t + 0.05,
        w - 0.2,
        0.35,
        title,
        font_size=title_size,
        bold=True,
        color=title_color,
    )
    lines = [("• " + b) for b in bullets]
    add_textbox_lines(
        slide,
        l + 0.1,
        t + 0.38,
        w - 0.2,
        h - 0.45,
        lines,
        font_size=bullet_size,
        color=bullet_color,
        line_spacing=1.2,
    )


def add_slide():
    return prs.slides.add_slide(blank_layout)


def footer(
    slide,
    text="Technical University of Munich  |  Chair of Urban Transport Systems",
    slide_num=None,
):
    add_rect(slide, 0, 7.22, 13.33, 0.28, fill_rgb=C_DARK_BLUE)
    add_textbox(
        slide,
        0.2,
        7.24,
        11.5,
        0.22,
        text,
        font_size=Pt(9),
        color=RGBColor(0xAA, 0xCC, 0xEE),
        align=PP_ALIGN.LEFT,
    )
    if slide_num is not None:
        add_textbox(
            slide,
            12.0,
            7.24,
            1.1,
            0.22,
            str(slide_num),
            font_size=Pt(9),
            bold=True,
            color=RGBColor(0xAA, 0xCC, 0xEE),
            align=PP_ALIGN.RIGHT,
        )


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=C_DARK_BLUE)
add_rect(s, 0, 2.6, 13.33, 2.4, fill_rgb=RGBColor(0x0E, 0x25, 0x40))
add_textbox(
    s,
    0.6,
    1.0,
    12.1,
    1.0,
    "ML Surrogates for Agent-Based Transport Models",
    font_size=Pt(32),
    bold=True,
    color=C_WHITE,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    s,
    0.6,
    2.0,
    12.1,
    0.5,
    "Uncertainty Quantification for GNN-Based Traffic Prediction",
    font_size=Pt(19),
    color=C_LIGHT_BLUE,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    s,
    0.6,
    2.75,
    12.1,
    0.45,
    "Thesis Defence Presentation",
    font_size=Pt(15),
    italic=True,
    color=RGBColor(0xAA, 0xCC, 0xEE),
    align=PP_ALIGN.CENTER,
)
add_textbox(
    s,
    0.6,
    3.3,
    12.1,
    0.45,
    "Supervisors: Dominik Natterer  |  Elena",
    font_size=Pt(14),
    color=C_LIGHT_BLUE,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    s,
    0.6,
    3.9,
    12.1,
    0.45,
    "Paris MATSim Dataset  •  1,000 Scenarios  •  31,635 Road Segments",
    font_size=Pt(13),
    color=RGBColor(0x88, 0xAA, 0xCC),
    align=PP_ALIGN.CENTER,
)
add_textbox(
    s,
    0.6,
    4.6,
    12.1,
    0.4,
    "March 2026",
    font_size=Pt(13),
    color=RGBColor(0x88, 0xAA, 0xCC),
    align=PP_ALIGN.CENTER,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "The Problem", "Why do we need an ML surrogate?", 2)

footer(s, slide_num=2)

add_figure(s, "fig11_thesis_workflow.png", 0.3, 1.25, 8.2, 5.0)
caption_box(
    s,
    0.3,
    6.25,
    8.2,
    "Fig: 1. Overall thesis pipeline: MATSim simulations → GNN surrogate → uncertainty-guided decisions",
)

# Right panel
add_rect(
    s,
    8.7,
    1.25,
    4.3,
    2.3,
    fill_rgb=RGBColor(0xFF, 0xF0, 0xE8),
    line_rgb=C_ORANGE,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    8.85,
    1.3,
    4.0,
    0.35,
    "Without Surrogate",
    font_size=Pt(13),
    bold=True,
    color=C_ORANGE,
)
add_textbox_lines(
    s,
    8.85,
    1.65,
    4.0,
    1.8,
    [
        "• 1 MATSim scenario = 8+ hours",
        "• Policy testing = weeks of compute",
        "• Very few alternatives evaluated",
    ],
    font_size=Pt(12),
    color=C_BLACK,
    line_spacing=1.25,
)

add_rect(
    s,
    8.7,
    3.7,
    4.3,
    2.3,
    fill_rgb=RGBColor(0xE8, 0xF5, 0xEC),
    line_rgb=C_GREEN,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    8.85,
    3.75,
    4.0,
    0.35,
    "With GNN Surrogate",
    font_size=Pt(13),
    bold=True,
    color=C_GREEN,
)
add_textbox_lines(
    s,
    8.85,
    4.1,
    4.0,
    1.8,
    [
        "• 1 scenario = seconds",
        "• BUT: needs uncertainty estimates",
        "• When can we trust the output?",
    ],
    font_size=Pt(12),
    color=C_BLACK,
    line_spacing=1.25,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Research Questions
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "4 Research Questions", "", 3)

footer(s, slide_num=3)


rqs = [
    (
        "RQ1",
        "MC Dropout Effectiveness",
        "How well does MC Dropout capture epistemic uncertainty in GNN traffic surrogates?\n(measured by Spearman ρ between predicted σ and actual error)",
    ),
    (
        "RQ2",
        "MC Dropout vs Ensembles",
        "How does MC Dropout compare to ensemble-based UQ in terms of uncertainty quality and computational efficiency?",
    ),
    (
        "RQ3",
        "Combining Identical Models",
        "Does combining uncertainty from architecturally identical models help? What role does diversity play?",
    ),
    (
        "RQ4",
        "Distribution-Free Coverage",
        "Can conformal prediction & temperature scaling transform raw uncertainty into trustworthy intervals with formal coverage guarantees?",
    ),
]
colors_bg = [
    RGBColor(0xE8, 0xF0, 0xFA),
    RGBColor(0xE8, 0xF5, 0xEC),
    RGBColor(0xFF, 0xF8, 0xE1),
    RGBColor(0xFC, 0xE8, 0xE8),
]
colors_brd = [C_MID_BLUE, C_GREEN, RGBColor(0xE0, 0xA0, 0x00), C_RED]
colors_hd = [C_MID_BLUE, C_GREEN, RGBColor(0xC0, 0x80, 0x00), C_RED]

for i, (rq_id, rq_title, rq_text) in enumerate(rqs):
    col = i % 2
    row = i // 2
    lx = 0.35 + col * 6.6
    ty = 1.25 + row * 2.85
    add_rect(
        s,
        lx,
        ty,
        6.3,
        2.6,
        fill_rgb=colors_bg[i],
        line_rgb=colors_brd[i],
        line_width=Pt(1.8),
    )
    add_textbox(
        s,
        lx + 0.15,
        ty + 0.1,
        1.1,
        0.4,
        rq_id,
        font_size=Pt(18),
        bold=True,
        color=colors_hd[i],
    )
    add_textbox(
        s,
        lx + 1.3,
        ty + 0.1,
        4.9,
        0.4,
        rq_title,
        font_size=Pt(14),
        bold=True,
        color=C_DARK_BLUE,
    )
    add_textbox(
        s,
        lx + 0.15,
        ty + 0.55,
        6.0,
        1.9,
        rq_text,
        font_size=Pt(12),
        color=C_BLACK,
        word_wrap=True,
    )

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset & Setup
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Dataset & Experimental Setup", "Paris MATSim Road Network", 4)

footer(s, slide_num=4)


add_figure(s, "fig_network_intro.png", 0.3, 1.2, 6.5, 5.3)
caption_box(
    s,
    0.3,
    6.5,
    6.5,
    "Fig: 2. Paris road network as directed graph. Coral nodes = capacity-reduced segments. "
    "Each node = one road segment. Target: Δv = volume change (veh/h).",
)

# Right: stats table
rows = [
    ("Road segments (nodes)", "31,635"),
    ("Network edges", "~85,000"),
    ("Scenarios used", "1,000 (10% of 10,000)"),
    ("Training scenarios", "800"),
    ("Validation scenarios", "100 (T7-T8)"),
    ("Test scenarios", "100 (T7-T8)"),
    ("Total test predictions", "3,163,500"),
    ("GPU", "NVIDIA T4 (Colab)"),
]
ty = 1.25
for label, val in rows:
    add_rect(s, 7.1, ty, 3.8, 0.36, fill_rgb=C_LIGHT_BLUE)
    add_rect(s, 10.9, ty, 2.15, 0.36, fill_rgb=RGBColor(0xD0, 0xE8, 0xFF))
    add_textbox(s, 7.2, ty + 0.04, 3.6, 0.3, label, font_size=Pt(11), color=C_DARK_BLUE)
    add_textbox(
        s,
        11.0,
        ty + 0.04,
        1.95,
        0.3,
        val,
        font_size=Pt(11),
        bold=True,
        color=C_DARK_BLUE,
        align=PP_ALIGN.CENTER,
    )
    ty += 0.42

add_textbox_lines(
    s,
    7.1,
    ty + 0.15,
    6.1,
    1.2,
    [
        ("5 Input Features per node:", True, C_DARK_BLUE),
        "• VOL_BASE_CASE  • CAPACITY_BASE_CASE",
        "• CAPACITY_REDUCTION  • FREESPEED  • LENGTH",
    ],
    font_size=Pt(12),
    line_spacing=1.3,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — GNN Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "GNN Architecture: PointNetTransfGAT",
    "6-layer graph neural network for node-level regression",
    5,
)

footer(s, slide_num=5)

add_figure(s, "fig8_architecture.png", 0.3, 1.2, 7.8, 5.2)
caption_box(
    s,
    0.3,
    6.4,
    7.8,
    "Fig: 3. PointNetTransfGAT: PointNetConv (geometry) → TransformerConv (attention) → GATConv (output). "
    "Input: 5 features per node. Output: Δv per node (veh/h).",
)

# Layer summary on right
layers = [
    ("Stage 1", "PointNetConv-1", "Start-point geometry encoding", "5 → 512"),
    ("Stage 2", "PointNetConv-2", "End-point geometry encoding", "512 → 128"),
    ("Stage 3", "TransformerConv-1", "4-head attention", "128 → 256"),
    ("Stage 4", "TransformerConv-2", "4-head attention", "256 → 512"),
    ("Stage 5", "GATConv", "Neighbourhood aggregation", "512 → 64"),
    ("Stage 6", "GATConv (final)", "Scalar prediction (T2-T8)", "64 → 1"),
]
ty = 1.25
for stage, name, desc, dims in layers:
    add_rect(s, 8.3, ty, 1.0, 0.56, fill_rgb=C_DARK_BLUE)
    add_textbox(
        s,
        8.3,
        ty + 0.1,
        1.0,
        0.35,
        stage,
        font_size=Pt(9),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_rect(s, 9.35, ty, 3.65, 0.56, fill_rgb=C_LIGHT_BLUE)
    add_textbox(
        s,
        9.45,
        ty + 0.02,
        3.5,
        0.25,
        name,
        font_size=Pt(11),
        bold=True,
        color=C_DARK_BLUE,
    )
    add_textbox(
        s,
        9.45,
        ty + 0.28,
        3.5,
        0.22,
        f"{desc}  [{dims}]",
        font_size=Pt(9),
        color=C_GRAY,
    )
    ty += 0.62

add_textbox(
    s,
    8.3,
    ty + 0.1,
    4.7,
    0.35,
    "Dropout applied in PointNetConv & TransformerConv layers",
    font_size=Pt(10),
    italic=True,
    color=C_ORANGE,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 8 Training Trials
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s, "8 Training Trials: Performance Summary", "Best model = T8 (R² = 0.5957)", 6
)

footer(s, slide_num=6)


add_figure(s, "fig1_trial_comparison.png", 0.3, 1.2, 7.5, 4.6)
caption_box(
    s,
    0.3,
    5.82,
    7.5,
    "Fig: 4. R², MAE, RMSE across Trials 2–8. T8 achieves best R²=0.5957, MAE=3.96 veh/h among UQ-compatible trials.",
)

# Table on right
headers = ["Trial", "R²", "MAE", "Dropout", "UQ?"]
data = [
    ("T1", "0.786", "2.97", "0.0", "❌"),
    ("T2", "0.512", "4.33", "0.3", "partial"),
    ("T3", "0.225", "5.99", "0.0", "❌"),
    ("T4", "0.243", "6.08", "0.0", "❌"),
    ("T5", "0.555", "4.24", "0.3", "✓"),
    ("T6", "0.522", "4.32", "0.3", "✓"),
    ("T7", "0.547", "4.06", "0.3", "✓"),
    ("T8★", "0.596", "3.96", "0.2", "✓ Best"),
]
col_w = [0.9, 0.75, 0.75, 0.9, 1.0]
col_x = [8.0, 8.95, 9.75, 10.55, 11.5]
ty = 1.2
# Header row
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(s, cx, ty, cw, 0.38, fill_rgb=C_DARK_BLUE)
    add_textbox(
        s,
        cx + 0.03,
        ty + 0.06,
        cw - 0.06,
        0.28,
        hdr,
        font_size=Pt(11),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
ty += 0.4
for row_i, row in enumerate(data):
    bg = (
        RGBColor(0xE0, 0xF0, 0xFF)
        if row[0] == "T8★"
        else (RGBColor(0xF0, 0xF0, 0xF0) if row_i % 2 == 0 else C_WHITE)
    )
    for ci, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(
            s,
            cx,
            ty,
            cw,
            0.38,
            fill_rgb=bg,
            line_rgb=RGBColor(0xCC, 0xCC, 0xCC),
            line_width=Pt(0.5),
        )
        fc = (
            C_ORANGE
            if row[0] == "T8★" and ci == 0
            else (C_GREEN if val == "✓ Best" else C_BLACK)
        )
        add_textbox(
            s,
            cx + 0.03,
            ty + 0.06,
            cw - 0.06,
            0.28,
            val,
            font_size=Pt(10),
            bold=(row[0] == "T8★"),
            color=fc,
            align=PP_ALIGN.CENTER,
        )
    ty += 0.4

add_textbox(
    s,
    8.0,
    ty + 0.05,
    4.7,
    0.45,
    "T1 excluded from UQ: zero dropout → σ=0 everywhere",
    font_size=Pt(10),
    italic=True,
    color=C_RED,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Why T1 Excluded & T9 Failed
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "Why T1 is Excluded & T9 Failed",
    "Understanding the boundaries of our UQ approach",
    7,
)

footer(s, slide_num=7)


add_figure(s, "fig12_trial_progression.png", 0.3, 1.2, 7.8, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    7.8,
    "Fig: 5. Trial progression T1→T8: R² (top) and MAE (bottom). T8 = best UQ-compatible trial.",
)

# T1 box
add_rect(
    s,
    8.3,
    1.2,
    4.7,
    2.45,
    fill_rgb=RGBColor(0xFF, 0xF0, 0xE8),
    line_rgb=C_RED,
    line_width=Pt(2),
)
add_textbox(
    s,
    8.45,
    1.25,
    4.4,
    0.4,
    "T1 — Excluded from UQ",
    font_size=Pt(14),
    bold=True,
    color=C_RED,
)
add_textbox_lines(
    s,
    8.45,
    1.68,
    4.4,
    1.9,
    [
        "• Has highest R² = 0.786 and MAE = 2.97",
        "• BUT: effective dropout = 0.0",
        "• All S=30 passes give identical output",
        "• σ = 0 everywhere → UQ undefined",
        "• Architectural constraint, not a flaw",
    ],
    font_size=Pt(11),
    color=C_BLACK,
    line_spacing=1.2,
)

# T9 box
add_rect(
    s,
    8.3,
    3.8,
    4.7,
    2.45,
    fill_rgb=RGBColor(0xFF, 0xEE, 0xEE),
    line_rgb=C_ORANGE,
    line_width=Pt(2),
)
add_textbox(
    s,
    8.45,
    3.85,
    4.4,
    0.4,
    "T9 — Heteroscedastic (Negative Result)",
    font_size=Pt(14),
    bold=True,
    color=C_ORANGE,
)
add_textbox_lines(
    s,
    8.45,
    4.28,
    4.4,
    1.9,
    [
        "• Tried: predict mean AND variance directly",
        "• Result: R² = 0.02 (target was ≥ 0.57)",
        "• Model inflated variance to absorb errors",
        "• Instead of improving mean predictions",
        "• Known failure mode (Seitzer et al.)",
    ],
    font_size=Pt(11),
    color=C_BLACK,
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 6 UQ Methods Overview
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s, "6 UQ Methods Evaluated", "All post-hoc — no model retraining required", 8
)

footer(s, slide_num=8)


methods = [
    (
        "1",
        "MC Dropout",
        "S=30 stochastic passes\nat inference time",
        C_MID_BLUE,
        RGBColor(0xD6, 0xE8, 0xF7),
    ),
    (
        "2",
        "Ensemble Variance",
        "5 seeded runs, std of\nper-run means",
        C_GREEN,
        RGBColor(0xE0, 0xF5, 0xE5),
    ),
    (
        "3",
        "Combined UQ",
        "√(σ²_MC + σ²_ens)\nquadrature sum",
        C_DARK_BLUE,
        C_LIGHT_BLUE,
    ),
    (
        "4",
        "Conformal Pred.",
        "Distribution-free\ncoverage guarantee",
        C_RED,
        RGBColor(0xFC, 0xE8, 0xE8),
    ),
    (
        "5",
        "Selective Pred.",
        "Filter by σ threshold,\nkeep confident ones",
        C_ORANGE,
        RGBColor(0xFF, 0xF3, 0xE0),
    ),
    (
        "6",
        "Temp. Scaling",
        "Post-hoc calibration\nwith scalar T",
        RGBColor(0x70, 0x30, 0xA0),
        RGBColor(0xF3, 0xE8, 0xFF),
    ),
]

positions = [(0.3, 1.2), (4.6, 1.2), (8.9, 1.2), (0.3, 4.1), (4.6, 4.1), (8.9, 4.1)]

for (lx, ty), (num, name, desc, brd, bg) in zip(positions, methods):
    add_rect(s, lx, ty, 4.0, 2.7, fill_rgb=bg, line_rgb=brd, line_width=Pt(2))
    add_rect(s, lx, ty, 0.55, 2.7, fill_rgb=brd)
    add_textbox(
        s,
        lx + 0.05,
        ty + 1.1,
        0.45,
        0.45,
        num,
        font_size=Pt(20),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        lx + 0.65,
        ty + 0.15,
        3.25,
        0.45,
        name,
        font_size=Pt(14),
        bold=True,
        color=brd,
    )
    add_textbox(
        s,
        lx + 0.65,
        ty + 0.65,
        3.25,
        1.9,
        desc,
        font_size=Pt(12),
        color=C_BLACK,
        word_wrap=True,
    )

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — RQ1: What is MC Dropout
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ1: What is MC Dropout?",
    "Keeping dropout active at inference to estimate uncertainty",
    9,
)

footer(s, slide_num=9)


add_figure(s, "fig13_mc_dropout_inference.png", 0.3, 1.2, 7.8, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    7.8,
    "Fig: 6. S=30 stochastic forward passes through T8. Per-node mean (ŷ) and std (σ) computed. "
    "σ serves as the uncertainty estimate. Spearman ρ=0.4820.",
)

# Algorithm steps on right
steps = [
    ("Step 1", "Load trained model with dropout layers"),
    ("Step 2", "Enable dropout at inference time"),
    ("Step 3", "Run S=30 forward passes on same graph"),
    ("Step 4", "Compute per-node mean: μ = (1/S)Σŷₛ"),
    ("Step 5", "Compute per-node std: σ = std(ŷ₁...ŷ₃₀)"),
    ("Step 6", "High σ = model uncertain about this road"),
]
ty = 1.25
for step, desc in steps:
    add_rect(s, 8.3, ty, 1.2, 0.52, fill_rgb=C_DARK_BLUE)
    add_textbox(
        s,
        8.3,
        ty + 0.1,
        1.2,
        0.32,
        step,
        font_size=Pt(10),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_rect(s, 9.55, ty, 3.45, 0.52, fill_rgb=C_LIGHT_BLUE)
    add_textbox(
        s, 9.65, ty + 0.08, 3.3, 0.38, desc, font_size=Pt(11), color=C_DARK_BLUE
    )
    ty += 0.58

add_rect(
    s,
    8.3,
    ty + 0.1,
    4.7,
    0.55,
    fill_rgb=RGBColor(0xFF, 0xF0, 0xD0),
    line_rgb=C_ORANGE,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    8.45,
    ty + 0.15,
    4.45,
    0.38,
    "S=30 chosen: plateau reached at S≈25, S=50 gives only +1% ρ",
    font_size=Pt(11),
    color=C_ORANGE,
    bold=True,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RQ1: MC Dropout Results
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ1 Results: MC Dropout Uncertainty Quality",
    "T8 achieves ρ = 0.4820 — best among all UQ-compatible trials",
    10,
)

footer(s, slide_num=10)


add_figure(s, "t8_per_graph_variation.png", 0.3, 1.2, 7.0, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    7.0,
    "Fig: 7. Per-graph Spearman ρ across all 100 test scenarios. "
    "Mean=0.464, Std=0.023 — consistent quality, no outlier scenarios.",
)

# Results table right
trial_rho = [
    ("T5", "0.3", "50", "0.4263"),
    ("T6", "0.3", "50", "0.4186"),
    ("T7", "0.3", "100", "0.4460"),
    ("T8★", "0.2", "100", "0.4820"),
]
headers2 = ["Trial", "Dropout", "Test\nGraphs", "Spearman ρ"]
col_w2 = [0.85, 0.9, 0.95, 1.05]
col_x2 = [7.6, 8.5, 9.45, 10.45]
ty = 1.2
for ci, (hdr, cx, cw) in enumerate(zip(headers2, col_x2, col_w2)):
    add_rect(s, cx, ty, cw, 0.45, fill_rgb=C_DARK_BLUE)
    add_textbox(
        s,
        cx + 0.03,
        ty + 0.08,
        cw - 0.06,
        0.35,
        hdr,
        font_size=Pt(10),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
ty += 0.47
for row_i, (trial, drop, graphs, rho) in enumerate(trial_rho):
    bg = (
        RGBColor(0xD6, 0xEA, 0xFF)
        if trial == "T8★"
        else (RGBColor(0xF5, 0xF5, 0xF5) if row_i % 2 == 0 else C_WHITE)
    )
    for ci, (val, cx, cw) in enumerate(zip([trial, drop, graphs, rho], col_x2, col_w2)):
        add_rect(
            s,
            cx,
            ty,
            cw,
            0.42,
            fill_rgb=bg,
            line_rgb=RGBColor(0xCC, 0xCC, 0xCC),
            line_width=Pt(0.5),
        )
        fc = C_ORANGE if trial == "T8★" and ci == 3 else C_BLACK
        add_textbox(
            s,
            cx + 0.03,
            ty + 0.08,
            cw - 0.06,
            0.28,
            val,
            font_size=Pt(11),
            bold=(trial == "T8★"),
            color=fc,
            align=PP_ALIGN.CENTER,
        )
    ty += 0.44

# Key metrics boxes
ty2 = ty + 0.2
metrics = [
    ("ρ = 0.4820", "Uncertainty-Error\nCorrelation", C_MID_BLUE),
    ("AUROC = 0.76", "Error Detection\n(top-10% errors)", C_GREEN),
    ("ΔR² = −0.010", "Accuracy Cost of\nEnabling Dropout", C_ORANGE),
]
for i, (val, label, col) in enumerate(metrics):
    lx = 7.5 + i * 1.95
    add_rect(
        s, lx, ty2, 1.85, 1.1, fill_rgb=C_LIGHT_BLUE, line_rgb=col, line_width=Pt(2)
    )
    add_textbox(
        s,
        lx + 0.05,
        ty2 + 0.08,
        1.75,
        0.45,
        val,
        font_size=Pt(16),
        bold=True,
        color=col,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        lx + 0.05,
        ty2 + 0.52,
        1.75,
        0.52,
        label,
        font_size=Pt(10),
        color=C_DARK_BLUE,
        align=PP_ALIGN.CENTER,
    )

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — S=30 Convergence
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "S = 30 Convergence Analysis",
    "Choosing the right number of MC forward passes",
    11,
)

footer(s, slide_num=11)


add_figure(s, "t8_s_convergence.png", 0.3, 1.2, 8.5, 4.9)
caption_box(
    s,
    0.3,
    6.15,
    8.5,
    "Fig: 8. Spearman ρ vs number of MC samples S (T8, 10 test graphs). "
    "Both ρ and mean σ plateau at S≈25. Circled point = S=30 (chosen operating point).",
)

add_textbox_lines(
    s,
    9.0,
    1.3,
    4.1,
    5.4,
    [
        ("Why S=30?", True, C_DARK_BLUE),
        "",
        ("Convergence plateau:", True, C_MID_BLUE),
        "• ρ flattens sharply after S≈25",
        "• S=30 → ρ = 0.4584",
        "• S=50 → ρ = 0.4632 (+1.0% only)",
        "",
        ("Compute cost:", True, C_MID_BLUE),
        "• S=30: 228 min on T4 GPU",
        "• S=50: ~380 min (+67% cost)",
        "• Gain: only +1% improvement",
        "",
        ("Conclusion:", True, C_GREEN),
        "S=30 is on the plateau.",
        "Optimal cost-quality trade-off.",
    ],
    font_size=Pt(12),
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — RQ2: Exp A MC Dropout vs Ensemble Variance
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ2: MC Dropout vs Ensemble Variance (Exp A)",
    "5 seeded runs on same T8 model",
    12,
)

footer(s, slide_num=12)


add_figure(s, "fig2_uq_ranking.png", 0.3, 1.2, 8.0, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    8.0,
    "Fig: 9. Spearman ρ for all evaluated UQ methods. "
    "MC Dropout (blue) outperforms ensemble variance (orange). Best = Exp A Combined (ρ=0.4909).",
)

# Results on right
results_a = [
    ("MC Dropout (S=30)", "0.4908", C_MID_BLUE, True),
    ("Ensemble Variance", "0.4370", C_ORANGE, False),
    ("Combined (MC+Ens)", "0.4909", C_GREEN, False),
]
ty = 1.3
add_textbox(
    s,
    8.5,
    ty,
    4.5,
    0.4,
    "Experiment A Results",
    font_size=Pt(15),
    bold=True,
    color=C_DARK_BLUE,
)
ty += 0.45
for method, rho, col, best in results_a:
    bg = RGBColor(0xE0, 0xF0, 0xFF) if best else RGBColor(0xF5, 0xF5, 0xF5)
    add_rect(s, 8.5, ty, 4.5, 0.65, fill_rgb=bg, line_rgb=col, line_width=Pt(1.5))
    add_textbox(
        s, 8.65, ty + 0.06, 3.0, 0.3, method, font_size=Pt(11), color=C_DARK_BLUE
    )
    add_textbox(
        s,
        11.65,
        ty + 0.06,
        1.2,
        0.3,
        f"ρ = {rho}",
        font_size=Pt(13),
        bold=True,
        color=col,
        align=PP_ALIGN.RIGHT,
    )
    ty += 0.72

add_rect(
    s,
    8.5,
    ty + 0.1,
    4.5,
    1.5,
    fill_rgb=RGBColor(0xE8, 0xF5, 0xEC),
    line_rgb=C_GREEN,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    8.65,
    ty + 0.18,
    4.25,
    0.35,
    "Key Finding",
    font_size=Pt(12),
    bold=True,
    color=C_GREEN,
)
add_textbox_lines(
    s,
    8.65,
    ty + 0.55,
    4.25,
    1.0,
    [
        "• MC Dropout beats Ensemble Variance by +12.3%",
        "• Combined adds negligible benefit (+0.01 ρ)",
        "• Ensemble variance is largely redundant",
    ],
    font_size=Pt(11),
    color=C_BLACK,
    line_spacing=1.2,
)

add_textbox(
    s,
    8.5,
    ty + 1.8,
    4.5,
    0.4,
    "Note: PyTorch Geometric API bug fixed before evaluation",
    font_size=Pt(10),
    italic=True,
    color=C_GRAY,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — RQ3: Multi-Model Ensemble (Exp B)
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ3: Multi-Model Ensemble (Exp B)",
    "T2 + T5 + T6 + T7 + T8 — weighted by individual R²",
    13,
)

footer(s, slide_num=13)


add_figure(s, "t7_vs_t8_uq_comparison.png", 0.3, 1.2, 7.8, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    7.8,
    "Fig: 10. Comprehensive UQ comparison: T7 vs T8 across AUROC, Spearman ρ, k₉₅, MAE, and selective prediction.",
)

# Exp B results
add_textbox(
    s,
    8.3,
    1.2,
    4.7,
    0.4,
    "Experiment B: 5-Model Ensemble",
    font_size=Pt(14),
    bold=True,
    color=C_DARK_BLUE,
)
metrics_b = [
    ("Ensemble ρ", "0.4333", C_ORANGE),
    ("Ensemble R²", "0.5656", C_ORANGE),
    ("Best single (T8) R²", "0.5957", C_GREEN),
    ("MC Dropout ρ (T8)", "0.4908", C_GREEN),
]
ty = 1.68
for label, val, col in metrics_b:
    add_rect(s, 8.3, ty, 3.0, 0.46, fill_rgb=C_LIGHT_BLUE)
    add_rect(s, 11.35, ty, 1.6, 0.46, fill_rgb=RGBColor(0xD8, 0xEE, 0xFF))
    add_textbox(
        s, 8.42, ty + 0.08, 2.85, 0.32, label, font_size=Pt(11), color=C_DARK_BLUE
    )
    add_textbox(
        s,
        11.4,
        ty + 0.08,
        1.5,
        0.32,
        val,
        font_size=Pt(12),
        bold=True,
        color=col,
        align=PP_ALIGN.CENTER,
    )
    ty += 0.52

add_rect(
    s,
    8.3,
    ty + 0.15,
    4.7,
    2.0,
    fill_rgb=RGBColor(0xFF, 0xF3, 0xE0),
    line_rgb=C_ORANGE,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    8.45,
    ty + 0.22,
    4.45,
    0.35,
    "Why ensemble doesn't help accuracy",
    font_size=Pt(12),
    bold=True,
    color=C_ORANGE,
)
add_textbox_lines(
    s,
    8.45,
    ty + 0.6,
    4.45,
    1.45,
    [
        "• Models span R² 0.51 (T2) to 0.60 (T8)",
        "• Averaging weaker models dilutes T8",
        "• Architectural diversity not tested here",
        "• UQ-wise: ρ=0.4333 < MC Dropout 0.4908",
    ],
    font_size=Pt(11),
    color=C_BLACK,
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — The PyG Bug (important finding)
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "Technical Finding: PyTorch Geometric API Bug",
    "Why ensemble results had to be corrected",
    14,
)

footer(s, slide_num=14)


add_rect(
    s,
    0.3,
    1.2,
    12.7,
    1.35,
    fill_rgb=RGBColor(0xFF, 0xEE, 0xDD),
    line_rgb=C_RED,
    line_width=Pt(2),
)
add_textbox(
    s,
    0.5,
    1.28,
    12.3,
    0.4,
    "Problem: Silent GATConv weight mismatch",
    font_size=Pt(16),
    bold=True,
    color=C_RED,
)
add_textbox(
    s,
    0.5,
    1.68,
    12.3,
    0.78,
    "Original ensemble scripts loaded checkpoints with strict=False. "
    "PyG 2.3.1 expected lin_src.weight / lin_dst.weight, but checkpoints stored lin.weight "
    "(older format). Weights silently dropped → R² dropped from ~0.59 to near-zero.",
    font_size=Pt(12),
    color=C_BLACK,
    word_wrap=True,
)

# Before/After
for i, (title, r2, rho, col, bg) in enumerate(
    [
        ("Before Fix (Corrupted)", "~0.0", "~0.0", C_RED, RGBColor(0xFF, 0xEE, 0xEE)),
        (
            "After Fix (Corrected)",
            "0.5656",
            "0.4333",
            C_GREEN,
            RGBColor(0xE8, 0xF5, 0xEC),
        ),
    ]
):
    lx = 0.5 + i * 6.3
    add_rect(s, lx, 2.75, 5.9, 2.5, fill_rgb=bg, line_rgb=col, line_width=Pt(2))
    add_textbox(
        s, lx + 0.15, 2.82, 5.6, 0.4, title, font_size=Pt(14), bold=True, color=col
    )
    add_textbox_lines(
        s,
        lx + 0.15,
        3.28,
        5.6,
        1.8,
        [f"Ensemble R²:   {r2}", f"Ensemble ρ:    {rho}"],
        font_size=Pt(20),
        bold=True,
        color=col,
        line_spacing=1.4,
    )

add_rect(
    s,
    0.3,
    5.45,
    12.7,
    0.9,
    fill_rgb=RGBColor(0xE8, 0xF0, 0xFF),
    line_rgb=C_MID_BLUE,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    0.5,
    5.52,
    12.3,
    0.35,
    "Fix: Remapped checkpoint keys before loading with strict=True",
    font_size=Pt(13),
    bold=True,
    color=C_MID_BLUE,
)
add_textbox(
    s,
    0.5,
    5.87,
    12.3,
    0.4,
    "Impact: MC Dropout, conformal prediction, temperature scaling results are NOT affected — "
    "they used the matching PyG version.",
    font_size=Pt(11),
    color=C_GRAY,
    word_wrap=True,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — RQ4: The Calibration Problem
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ4: The Calibration Problem",
    "Raw MC Dropout σ is NOT a calibrated standard deviation",
    15,
)

footer(s, slide_num=15)


add_figure(s, "fig7_calibration.png", 0.3, 1.2, 5.5, 4.5)
caption_box(
    s,
    0.3,
    5.72,
    5.5,
    "Fig: 11. Scaling factor k₉₅ needed for ±kσ to achieve 95% coverage. "
    "Ideal Gaussian: k=1.96. T8 MC Dropout requires k=11.34 — nearly 6× larger.",
)

add_figure(s, "t8_reliability_diagram.png", 6.0, 1.2, 7.0, 4.5)
caption_box(
    s,
    6.0,
    5.72,
    7.0,
    "Fig: 12. Reliability diagram (left): observed coverage always below diagonal = systematic under-coverage. "
    "ECE = 0.265 means coverage deviates 26.5 pp on average.",
)

# Bottom summary strip
add_rect(
    s,
    0.3,
    6.2,
    12.7,
    0.95,
    fill_rgb=RGBColor(0xFF, 0xEE, 0xDD),
    line_rgb=C_ORANGE,
    line_width=Pt(1.5),
)
add_textbox_lines(
    s,
    0.5,
    6.27,
    12.3,
    0.8,
    [
        (
            "Key insight: σ is a ranking signal — NOT a probability. "
            "Spearman ρ=0.4820 works perfectly. "
            "But ±1.96σ only gives 54.8% coverage (target: 95%).",
            False,
            C_ORANGE,
        )
    ],
    font_size=Pt(12),
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Coverage Comparison Table
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ4: Coverage — MC Dropout vs Conformal",
    "Conformal prediction fixes the coverage gap",
    16,
)

footer(s, slide_num=16)


add_figure(s, "t8_calibration_curve.png", 0.3, 1.2, 6.8, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    6.8,
    "Fig: 13. T8 calibration curve. Raw MC Dropout (blue) severely undercovers. "
    "Global conformal (orange) and adaptive conformal (green) both hit nominal targets.",
)

# Coverage table on right
add_textbox(
    s,
    7.4,
    1.2,
    5.6,
    0.38,
    "Coverage Comparison Table",
    font_size=Pt(14),
    bold=True,
    color=C_DARK_BLUE,
)
coverage_data = [
    ("50%", "23.3%", "50.0%"),
    ("80%", "40.1%", "80.0%"),
    ("90%", "48.6%", "90.0%"),
    ("95%", "54.8%", "95.0%"),
]
hdrs3 = ["Nominal", "Raw MC\nCoverage", "Conformal\nCoverage"]
col_w3 = [1.05, 1.4, 1.6]
col_x3 = [7.4, 8.5, 9.95]
ty = 1.65
for ci, (hdr, cx, cw) in enumerate(zip(hdrs3, col_x3, col_w3)):
    add_rect(s, cx, ty, cw, 0.5, fill_rgb=C_DARK_BLUE)
    add_textbox(
        s,
        cx + 0.03,
        ty + 0.06,
        cw - 0.06,
        0.4,
        hdr,
        font_size=Pt(11),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
ty += 0.52
for nominal, mc_cov, conf_cov in coverage_data:
    for ci, (val, cx, cw) in enumerate(
        zip([nominal, mc_cov, conf_cov], col_x3, col_w3)
    ):
        bg = (
            RGBColor(0xFF, 0xEE, 0xDD)
            if ci == 1
            else (RGBColor(0xE8, 0xF5, 0xEC) if ci == 2 else C_LIGHT_BLUE)
        )
        fc = C_RED if ci == 1 else (C_GREEN if ci == 2 else C_DARK_BLUE)
        add_rect(
            s,
            cx,
            ty,
            cw,
            0.48,
            fill_rgb=bg,
            line_rgb=RGBColor(0xCC, 0xCC, 0xCC),
            line_width=Pt(0.5),
        )
        add_textbox(
            s,
            cx + 0.03,
            ty + 0.1,
            cw - 0.06,
            0.3,
            val,
            font_size=Pt(13),
            bold=(ci > 0),
            color=fc,
            align=PP_ALIGN.CENTER,
        )
    ty += 0.5

add_rect(
    s,
    7.4,
    ty + 0.12,
    4.25,
    1.8,
    fill_rgb=RGBColor(0xE8, 0xF5, 0xEC),
    line_rgb=C_GREEN,
    line_width=Pt(1.8),
)
add_textbox(
    s,
    7.55,
    ty + 0.18,
    3.95,
    0.4,
    "Why conformal works",
    font_size=Pt(12),
    bold=True,
    color=C_GREEN,
)
add_textbox_lines(
    s,
    7.55,
    ty + 0.6,
    3.95,
    1.2,
    [
        "• Uses empirical calibration quantile",
        "• Distribution-free — no assumptions",
        "• Guaranteed by construction",
        "• Cost: wider fixed intervals",
    ],
    font_size=Pt(11),
    color=C_BLACK,
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Temperature Scaling
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ4: Temperature Scaling — Post-Hoc Calibration",
    "Single scalar T=2.70 reduces ECE by 82%",
    17,
)

footer(s, slide_num=17)


add_figure(s, "t8_temperature_scaling.png", 0.3, 1.2, 8.2, 5.1)
caption_box(
    s,
    0.3,
    6.35,
    8.2,
    "Fig: 14. Top: reliability diagrams before (ECE=0.269) and after (ECE=0.048) scaling with T=2.70. "
    "Bottom: coverage gap reduction at each nominal level. Evaluated on 80 held-out graphs.",
)

# Right panel
add_textbox(
    s,
    8.75,
    1.25,
    4.3,
    0.4,
    "What is Temperature Scaling?",
    font_size=Pt(13),
    bold=True,
    color=C_DARK_BLUE,
)
add_textbox(
    s,
    8.75,
    1.68,
    4.3,
    0.7,
    "σ_scaled = σ_raw × T\n"
    "Optimal T learned on calibration set (20 graphs) by minimising ECE.",
    font_size=Pt(11),
    color=C_BLACK,
    word_wrap=True,
)

metrics_ts = [
    ("T = 2.70", "Optimal scaling factor", C_MID_BLUE),
    ("ECE: 0.265 → 0.048", "82% improvement", C_GREEN),
    ("NLL: 21.65 → 4.75", "78% improvement", C_GREEN),
    ("KS stat: 0.245 → 0.104", "PIT improvement", C_GREEN),
    ("95% coverage: 83.3%", "Not 95% — residual gap", C_ORANGE),
]
ty = 2.55
for val, desc, col in metrics_ts:
    add_rect(s, 8.75, ty, 4.3, 0.5, fill_rgb=C_LIGHT_BLUE)
    add_textbox(
        s, 8.9, ty + 0.06, 2.8, 0.35, val, font_size=Pt(12), bold=True, color=col
    )
    add_textbox(
        s,
        11.75,
        ty + 0.06,
        1.2,
        0.35,
        desc,
        font_size=Pt(9),
        italic=True,
        color=C_GRAY,
        align=PP_ALIGN.RIGHT,
    )
    ty += 0.56

add_rect(
    s,
    8.75,
    ty + 0.05,
    4.3,
    0.7,
    fill_rgb=RGBColor(0xFF, 0xF3, 0xE0),
    line_rgb=C_ORANGE,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    8.9,
    ty + 0.1,
    4.1,
    0.55,
    "Useful middle step. Residual gap at tails → conformal prediction still needed for guarantees.",
    font_size=Pt(11),
    color=C_ORANGE,
    word_wrap=True,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Conformal Prediction — How It Works
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ4: Split Conformal Prediction",
    "Distribution-free coverage guarantee — no model assumptions needed",
    18,
)

footer(s, slide_num=18)


add_figure(s, "fig14_conformal_workflow.png", 0.3, 1.2, 8.3, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    8.3,
    "Fig: 15. Conformal prediction workflow for T8. 100 test graphs → 50 calibration + 50 evaluation. "
    "Calibration residuals define quantile q. Intervals ŷ±q guarantee ≥90% or ≥95% coverage.",
)

# Right: algorithm steps
add_textbox(
    s,
    8.85,
    1.25,
    4.2,
    0.4,
    "How It Works",
    font_size=Pt(14),
    bold=True,
    color=C_DARK_BLUE,
)
steps_cp = [
    ("1", "Split test data: 50 calibration + 50 evaluation graphs"),
    ("2", "On calibration set: compute residuals |y − ŷ| per node"),
    ("3", "Find quantile q at desired level (e.g. 90th or 95th)"),
    ("4", "Apply to eval set: interval = [ŷ − q, ŷ + q]"),
    ("5", "Coverage guaranteed: ≥90% or ≥95% by construction"),
]
ty = 1.72
for num, desc in steps_cp:
    add_rect(s, 8.85, ty, 0.42, 0.5, fill_rgb=C_MID_BLUE)
    add_textbox(
        s,
        8.85,
        ty + 0.08,
        0.42,
        0.35,
        num,
        font_size=Pt(13),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_rect(s, 9.32, ty, 3.7, 0.5, fill_rgb=C_LIGHT_BLUE)
    add_textbox(
        s,
        9.42,
        ty + 0.08,
        3.55,
        0.38,
        desc,
        font_size=Pt(11),
        color=C_DARK_BLUE,
        word_wrap=True,
    )
    ty += 0.56

add_rect(
    s,
    8.85,
    ty + 0.1,
    4.18,
    0.85,
    fill_rgb=RGBColor(0xE8, 0xF5, 0xEC),
    line_rgb=C_GREEN,
    line_width=Pt(2),
)
add_textbox_lines(
    s,
    9.0,
    ty + 0.17,
    3.9,
    0.65,
    [
        ("Key: Guarantee holds regardless of model internals", True, C_GREEN),
        "Only assumption: data is exchangeable",
    ],
    font_size=Pt(11),
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Conformal Results
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ4: Conformal Prediction Results",
    "Near-exact coverage achieved at both levels",
    19,
)

footer(s, slide_num=19)


add_figure(s, "fig3_conformal_coverage.png", 0.3, 1.2, 7.2, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    7.2,
    "Fig: 16. Nominal vs achieved coverage (left) and interval half-widths (right). "
    "Both 90% and 95% levels meet targets within 0.02 percentage points. "
    "Evaluated on 50 graphs (1,581,750 nodes).",
)

# Results boxes right
add_textbox(
    s,
    7.75,
    1.25,
    5.3,
    0.4,
    "Coverage Results (T8)",
    font_size=Pt(15),
    bold=True,
    color=C_DARK_BLUE,
)
conf_results = [
    ("90% level", "q = 9.92 veh/h", "Achieved: 90.02%", C_MID_BLUE),
    ("95% level", "q = 14.68 veh/h", "Achieved: 95.01%", C_GREEN),
]
ty = 1.72
for level, quantile, achieved, col in conf_results:
    add_rect(
        s, 7.75, ty, 5.3, 1.3, fill_rgb=C_LIGHT_BLUE, line_rgb=col, line_width=Pt(2)
    )
    add_textbox(
        s, 7.9, ty + 0.08, 5.0, 0.38, level, font_size=Pt(14), bold=True, color=col
    )
    add_textbox(
        s, 7.9, ty + 0.46, 5.0, 0.3, quantile, font_size=Pt(12), color=C_DARK_BLUE
    )
    add_textbox(
        s, 7.9, ty + 0.78, 5.0, 0.38, achieved, font_size=Pt(15), bold=True, color=col
    )
    ty += 1.4

add_rect(
    s,
    7.75,
    ty + 0.05,
    5.3,
    2.25,
    fill_rgb=RGBColor(0xF0, 0xF8, 0xE8),
    line_rgb=C_GREEN,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    7.9,
    ty + 0.12,
    5.05,
    0.35,
    "What this means for planners",
    font_size=Pt(12),
    bold=True,
    color=C_GREEN,
)
add_textbox_lines(
    s,
    7.9,
    ty + 0.5,
    5.05,
    1.65,
    [
        "• At least 95% of true Δv values fall within ±14.68 veh/h",
        "• This guarantee is mathematically proven",
        "• Fixed-width: same interval for every road segment",
        "• Cost: ~6× wider than raw MC Dropout interval",
    ],
    font_size=Pt(11),
    color=C_BLACK,
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Adaptive Conformal Prediction
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "RQ4: Adaptive Conformal Prediction",
    "Combining MC Dropout ranking with conformal guarantees",
    20,
)

footer(s, slide_num=20)


add_figure(s, "t8_conformal_conditional.png", 0.3, 1.2, 8.5, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    8.5,
    "Fig: 17. Left: standard conformal over-covers low-uncertainty nodes (98.6%) and under-covers high-uncertainty (62.9%). "
    "Right: adaptive conformal narrows this to [90.0%, 96.2%] — much more uniform conditional coverage.",
)

# Right panel
add_textbox(
    s,
    9.05,
    1.25,
    4.0,
    0.4,
    "Standard vs Adaptive",
    font_size=Pt(14),
    bold=True,
    color=C_DARK_BLUE,
)
add_textbox_lines(
    s,
    9.05,
    1.72,
    4.0,
    1.5,
    [
        ("Standard conformal:", True, C_ORANGE),
        "• Fixed width ±q for all nodes",
        "• D1 (low σ): 98.6% coverage",
        "• D10 (high σ): 62.9% coverage",
    ],
    font_size=Pt(12),
    line_spacing=1.2,
)
add_textbox_lines(
    s,
    9.05,
    3.35,
    4.0,
    1.5,
    [
        ("Adaptive conformal:", True, C_GREEN),
        "• Width = q_adapt × σ (node-specific)",
        "• Wider where model is uncertain",
        "• Narrows disparity → [90%, 96.2%]",
    ],
    font_size=Pt(12),
    line_spacing=1.2,
)
add_rect(
    s,
    9.05,
    4.98,
    4.0,
    1.2,
    fill_rgb=C_LIGHT_BLUE,
    line_rgb=C_MID_BLUE,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    9.2,
    5.04,
    3.75,
    0.35,
    "Best of both worlds:",
    font_size=Pt(12),
    bold=True,
    color=C_MID_BLUE,
)
add_textbox(
    s,
    9.2,
    5.4,
    3.75,
    0.7,
    "MC Dropout σ provides adaptive widths.\nConformal guarantees marginal coverage.",
    font_size=Pt(11),
    color=C_DARK_BLUE,
    word_wrap=True,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Selective Prediction
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "Selective Prediction: Filter by Uncertainty",
    "Keep only confident predictions → dramatically lower MAE",
    21,
)

footer(s, slide_num=21)


add_figure(s, "t8_selective_prediction_curve.png", 0.3, 1.2, 7.5, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    7.5,
    "Fig: 18. MAE as fraction of retained predictions (sorted by ascending σ) for T8. "
    "Retaining the 50% most certain predictions gives −41.2% MAE. "
    "Based on 100 test graphs (3,163,500 nodes).",
)

# Table on right
sel_data = [
    ("100%", "3.95", "—", C_RED),
    ("90%", "3.23", "−18.3%", C_ORANGE),
    ("50%", "2.32", "−41.2%", C_MID_BLUE),
    ("25%", "1.79", "−54.6%", C_GREEN),
    ("10%", "1.06", "−73.3%", C_GREEN),
]
hdrs_sel = ["Retained", "MAE (veh/h)", "Reduction"]
col_w_sel = [1.2, 1.4, 1.3]
col_x_sel = [8.05, 9.3, 10.75]
ty = 1.2
for hdr, cx, cw in zip(hdrs_sel, col_x_sel, col_w_sel):
    add_rect(s, cx, ty, cw, 0.45, fill_rgb=C_DARK_BLUE)
    add_textbox(
        s,
        cx + 0.03,
        ty + 0.07,
        cw - 0.06,
        0.32,
        hdr,
        font_size=Pt(11),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
ty += 0.47
for ret, mae, red, col in sel_data:
    highlight = ret == "50%"
    bg = RGBColor(0xD0, 0xE8, 0xFF) if highlight else RGBColor(0xF5, 0xF5, 0xF5)
    for val, cx, cw in zip([ret, mae, red], col_x_sel, col_w_sel):
        add_rect(
            s,
            cx,
            ty,
            cw,
            0.48,
            fill_rgb=bg,
            line_rgb=col if highlight else RGBColor(0xCC, 0xCC, 0xCC),
            line_width=Pt(1.5 if highlight else 0.5),
        )
        add_textbox(
            s,
            cx + 0.03,
            ty + 0.09,
            cw - 0.06,
            0.32,
            val,
            font_size=Pt(12),
            bold=highlight,
            color=col,
            align=PP_ALIGN.CENTER,
        )
    ty += 0.5

add_rect(
    s,
    8.05,
    ty + 0.15,
    4.0,
    1.15,
    fill_rgb=RGBColor(0xE8, 0xF5, 0xEC),
    line_rgb=C_GREEN,
    line_width=Pt(1.5),
)
add_textbox(
    s,
    8.2,
    ty + 0.22,
    3.75,
    0.35,
    "Practical use for planners",
    font_size=Pt(12),
    bold=True,
    color=C_GREEN,
)
add_textbox_lines(
    s,
    8.2,
    ty + 0.6,
    3.75,
    0.65,
    [
        "Accept bottom 50% σ → expect MAE ≈ 2.32 veh/h",
        "Run full MATSim only for top 10% most uncertain",
    ],
    font_size=Pt(11),
    color=C_BLACK,
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — 3-Layer UQ Hierarchy
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s, "The 3-Layer UQ Hierarchy", "Main practical recommendation of this thesis", 22
)

footer(s, slide_num=22)


layers_info = [
    (
        "LAYER 1",
        "MC Dropout (S=30)",
        "Purpose: RANKING predictions by reliability",
        [
            "• σ = per-node uncertainty estimate",
            "• Spearman ρ = 0.4820 with actual error",
            "• AUROC = 0.76 for error detection",
            "• Use: selective prediction, flagging",
        ],
        C_MID_BLUE,
        RGBColor(0xD6, 0xE8, 0xF7),
    ),
    (
        "LAYER 2",
        "Temperature Scaling (T=2.70)",
        "Purpose: CALIBRATE σ toward probability",
        [
            "• σ_scaled = σ_raw × 2.70",
            "• ECE reduced 82% (0.265 → 0.048)",
            "• Better Gaussian approximation",
            "• Use: when approximate probabilities needed",
        ],
        RGBColor(0x70, 0x30, 0xA0),
        RGBColor(0xF0, 0xE8, 0xFF),
    ),
    (
        "LAYER 3",
        "Conformal Prediction",
        "Purpose: GUARANTEE coverage bounds",
        [
            "• 90.02% at 90% nominal ✓",
            "• 95.01% at 95% nominal ✓",
            "• Distribution-free — no assumptions",
            "• Use: contractual / operational requirements",
        ],
        C_GREEN,
        RGBColor(0xE8, 0xF5, 0xEC),
    ),
]

for i, (lyr, name, purpose, bullets, col, bg) in enumerate(layers_info):
    lx = 0.35 + i * 4.35
    add_rect(s, lx, 1.2, 4.1, 0.45, fill_rgb=col)
    add_textbox(
        s, lx + 0.1, 1.25, 3.9, 0.35, lyr, font_size=Pt(13), bold=True, color=C_WHITE
    )
    add_rect(s, lx, 1.65, 4.1, 5.1, fill_rgb=bg, line_rgb=col, line_width=Pt(2))
    add_textbox(
        s, lx + 0.15, 1.72, 3.85, 0.42, name, font_size=Pt(14), bold=True, color=col
    )
    add_textbox(
        s,
        lx + 0.15,
        2.15,
        3.85,
        0.35,
        purpose,
        font_size=Pt(11),
        italic=True,
        color=C_DARK_BLUE,
    )
    add_textbox_lines(
        s,
        lx + 0.15,
        2.6,
        3.85,
        4.0,
        bullets,
        font_size=Pt(11),
        color=C_BLACK,
        line_spacing=1.25,
    )

# Arrow between layers
for i in range(2):
    lx_arrow = 4.45 + i * 4.35 - 0.02
    add_textbox(
        s,
        lx_arrow,
        3.6,
        0.3,
        0.5,
        "→",
        font_size=Pt(28),
        bold=True,
        color=C_GRAY,
        align=PP_ALIGN.CENTER,
    )

add_rect(s, 0.35, 6.75, 12.63, 0.4, fill_rgb=C_DARK_BLUE)
add_textbox(
    s,
    0.5,
    6.8,
    12.4,
    0.28,
    "All three methods are post-hoc: no model retraining required. "
    "They can be applied to any trained GNN.",
    font_size=Pt(11),
    color=C_LIGHT_BLUE,
    align=PP_ALIGN.CENTER,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — UQ Methods Summary Table
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "UQ Methods: Summary Comparison",
    "All 5 methods evaluated on T8 (100 test graphs, 3.16M predictions)",
    23,
)

footer(s, slide_num=23)


summary_rows = [
    (
        "MC Dropout\n(S=30)",
        "ρ=0.4820\nCRPS/MAE=0.857",
        "Node-specific ranking\nSelective pred. −41.2% MAE",
        "Not calibrated (k₉₅=11.34)\nS-fold compute cost",
        C_MID_BLUE,
    ),
    (
        "Temp. Scaling\n(T=2.70)",
        "ECE: 0.265→0.048\nKS: 0.245→0.104",
        "Single parameter\nNo retraining",
        "Gaussian assumption\nResidual miscal. at tails",
        RGBColor(0x70, 0x30, 0xA0),
    ),
    (
        "Standard\nConformal",
        "PICP 90%: 90.02%\nPICP 95%: 95.01%",
        "Coverage guarantee\nDistribution-free",
        "Fixed-width intervals\nPoor conditional coverage",
        C_GREEN,
    ),
    (
        "Adaptive\nConformal",
        "Cond. cov.\n[90.0%, 96.2%]",
        "Node-specific widths\nConditional coverage",
        "Requires σ from MC Dropout",
        C_GREEN,
    ),
    (
        "Multi-Model\nEnsemble",
        "ρ=0.4333\nR²=0.5656",
        "Complementary signal\nNo MC overhead",
        "Weaker than MC Dropout\nNo accuracy gain over T8 alone",
        C_ORANGE,
    ),
]
col_labels = ["Method", "Key Metric", "Strength", "Limitation"]
col_widths = [1.8, 2.8, 3.2, 3.4]
col_starts = [0.3, 2.15, 5.0, 8.25]

ty = 1.2
for hdr, cx, cw in zip(col_labels, col_starts, col_widths):
    add_rect(s, cx, ty, cw, 0.45, fill_rgb=C_DARK_BLUE)
    add_textbox(
        s,
        cx + 0.05,
        ty + 0.08,
        cw - 0.1,
        0.3,
        hdr,
        font_size=Pt(12),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
ty += 0.47
for method, metric, strength, limitation, col in summary_rows:
    row_h = 0.98
    for val, cx, cw in zip(
        [method, metric, strength, limitation], col_starts, col_widths
    ):
        add_rect(
            s,
            cx,
            ty,
            cw,
            row_h,
            fill_rgb=C_LIGHT_BLUE,
            line_rgb=RGBColor(0xBB, 0xCC, 0xDD),
            line_width=Pt(0.5),
        )
        add_textbox(
            s,
            cx + 0.07,
            ty + 0.08,
            cw - 0.14,
            row_h - 0.12,
            val,
            font_size=Pt(10),
            color=C_BLACK,
            word_wrap=True,
        )
    # colour the method cell
    add_rect(
        s,
        col_starts[0],
        ty,
        col_widths[0],
        row_h,
        fill_rgb=C_LIGHT_BLUE,
        line_rgb=col,
        line_width=Pt(2),
    )
    add_textbox(
        s,
        col_starts[0] + 0.07,
        ty + 0.08,
        col_widths[0] - 0.14,
        row_h - 0.12,
        method,
        font_size=Pt(11),
        bold=True,
        color=col,
        word_wrap=True,
    )
    ty += row_h + 0.06

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — What is Being Predicted (Traffic context)
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "Traffic Context: What is Being Predicted?",
    "Δv = change in traffic volume per road segment after policy intervention",
    24,
)

footer(s, slide_num=24)


add_figure(s, "fig3_feature_distributions.png", 0.3, 1.2, 7.5, 4.8)
caption_box(
    s,
    0.3,
    6.05,
    7.5,
    "Fig: 19. Distribution of 5 input features across 200 training graphs (6.33M node observations). "
    "VOL_BASE_CASE and CAPACITY_BASE_CASE are right-skewed (heavy-tailed urban traffic). "
    "FREESPEED shows discrete speed limits. LENGTH spans several orders of magnitude.",
)

add_textbox_lines(
    s,
    8.05,
    1.25,
    5.0,
    5.5,
    [
        ("What is Δv?", True, C_DARK_BLUE),
        "The change in traffic volume on each",
        "road segment after a policy applies:",
        "",
        "  Δv = v_policy − v_baseline  (veh/h)",
        "",
        ("Policy = capacity reduction:", True, C_DARK_BLUE),
        "• 10%–100% of road capacity reduced",
        "• Applied to a random subset of roads",
        "• Simulates: lane closures, traffic calming,",
        "  bus-only lanes, congestion pricing...",
        "",
        ("Why hard to predict?", True, C_ORANGE),
        "• Traffic redistributes across network",
        "• High-flow roads change most but vary most",
        "• ρ ≈ 0.74–0.81 (low-flow) vs 0.26–0.41 (high-flow)",
    ],
    font_size=Pt(12),
    line_spacing=1.2,
)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — 3-Tier Decision Framework
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "Traffic Policy: 3-Tier Decision Framework",
    "Uncertainty-guided workflow for urban planners",
    25,
)

footer(s, slide_num=25)


add_figure(s, "fig9_policy_explanation.png", 0.3, 1.2, 7.8, 5.1)
caption_box(
    s,
    0.3,
    6.35,
    7.8,
    "Fig: 20. Uncertainty-guided decision framework. GNN surrogate (T8, S=30) routes each prediction "
    "based on MC Dropout σ. Bottom 50% = ACCEPT, 50–90% = FLAG, Top 10% = REJECT for re-simulation.",
)

# 3 tier boxes
tiers = [
    (
        "ACCEPT",
        "Bottom 50% of σ",
        "MAE ≈ 2.32 veh/h\n−41.2% error",
        "Use surrogate output directly\nfor policy decision",
        C_GREEN,
    ),
    (
        "FLAG",
        "50th–90th percentile",
        "Some uncertainty\npresent",
        "Sensitivity analysis\nor cross-check needed",
        C_ORANGE,
    ),
    (
        "REJECT",
        "Top 10% of σ",
        "High uncertainty\nrun full MATSim",
        "Route to full 8-hr\nMATSim simulation",
        C_RED,
    ),
]
ty = 1.25
for tier, sigma_range, detail, action, col in tiers:
    add_rect(s, 8.3, ty, 1.2, 1.65, fill_rgb=col)
    add_textbox(
        s,
        8.3,
        ty + 0.55,
        1.2,
        0.5,
        tier,
        font_size=Pt(13),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_rect(s, 9.55, ty, 1.9, 1.65, fill_rgb=C_LIGHT_BLUE)
    add_textbox(
        s,
        9.65,
        ty + 0.1,
        1.75,
        0.35,
        sigma_range,
        font_size=Pt(11),
        bold=True,
        color=col,
    )
    add_textbox(
        s,
        9.65,
        ty + 0.48,
        1.75,
        1.0,
        detail,
        font_size=Pt(10),
        color=C_DARK_BLUE,
        word_wrap=True,
    )
    add_rect(s, 11.5, ty, 1.55, 1.65, fill_rgb=RGBColor(0xF0, 0xF0, 0xF0))
    add_textbox(
        s,
        11.6,
        ty + 0.38,
        1.35,
        0.9,
        action,
        font_size=Pt(10),
        color=C_DARK_BLUE,
        word_wrap=True,
    )
    ty += 1.78

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Impact for Urban Planning
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(
    s,
    "Impact for Urban Planning",
    "How this thesis changes the policy evaluation workflow",
    26,
)

footer(s, slide_num=26)


add_figure(s, "t8_error_detection_auroc.png", 0.3, 1.2, 6.8, 4.5)
caption_box(
    s,
    0.3,
    5.72,
    6.8,
    "Fig: 21. ROC curves for T8 MC Dropout as error detector. "
    "AUROC=0.76 (top-10% errors) and 0.74 (top-20% errors) — well above 0.5 random baseline. "
    "Model can identify unreliable predictions before a planner acts on them.",
)

impact_pts = [
    (
        C_GREEN,
        "Faster policy screening",
        "Seconds per scenario instead of 8+ hours. "
        "Hundreds of alternatives evaluated in hours.",
    ),
    (
        C_MID_BLUE,
        "Reliable uncertainty scores",
        "ρ=0.4820: high-σ roads reliably indicate where model may be wrong. "
        "AUROC=0.76 for flagging top-10% errors.",
    ),
    (
        C_ORANGE,
        "Smart resource allocation",
        "Expensive MATSim re-simulation only for top-10% uncertain predictions. "
        "90% of predictions can be accepted or flagged.",
    ),
    (
        C_RED,
        "Statistically valid intervals",
        "Conformal prediction gives ±14.68 veh/h interval with 95% coverage guarantee "
        "— suitable for formal planning documents.",
    ),
]
ty = 1.25
for col, title, text in impact_pts:
    add_rect(s, 7.3, ty, 0.25, 0.9, fill_rgb=col)
    add_rect(s, 7.6, ty, 5.4, 0.9, fill_rgb=C_LIGHT_BLUE)
    add_textbox(
        s, 7.72, ty + 0.05, 5.2, 0.32, title, font_size=Pt(13), bold=True, color=col
    )
    add_textbox(
        s,
        7.72,
        ty + 0.37,
        5.2,
        0.48,
        text,
        font_size=Pt(11),
        color=C_DARK_BLUE,
        word_wrap=True,
    )
    ty += 1.05

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Answers to All 4 RQs
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Summary: Answers to All 4 Research Questions", "", 27)

footer(s, slide_num=27)


rq_answers = [
    (
        "RQ1",
        "MC Dropout Effectiveness",
        "T8 achieves Spearman ρ = 0.4820 (100 test graphs, 3.16M predictions). "
        "Consistent per-graph ρ = 0.464 ± 0.023. AUROC = 0.76 for top-10% error detection. "
        "Enabling dropout costs only ΔR² = −0.010.",
        C_MID_BLUE,
        RGBColor(0xD6, 0xE8, 0xF7),
    ),
    (
        "RQ2",
        "MC Dropout vs Ensembles",
        "MC Dropout (ρ=0.4908) outperforms ensemble variance (ρ=0.4370) by +12.3% "
        "and multi-model disagreement (ρ=0.4333) by +13.3%. "
        "MC Dropout is best and most computationally efficient single method.",
        C_GREEN,
        RGBColor(0xE0, 0xF5, 0xE5),
    ),
    (
        "RQ3",
        "Combining Identical Models",
        "Combined MC+Ensemble (ρ=0.4909) barely beats MC alone (ρ=0.4908): negligible gain. "
        "Multi-model ensemble R²=0.5656 < best single T8 R²=0.5957. "
        "Architectural diversity not tested — may unlock further benefits.",
        RGBColor(0xC0, 0x80, 0x00),
        RGBColor(0xFF, 0xF8, 0xE0),
    ),
    (
        "RQ4",
        "Distribution-Free Coverage",
        "Yes. Conformal: 90.02% @ 90%, 95.01% @ 95% — guaranteed. "
        "Temperature scaling reduces ECE by 82%. "
        "Adaptive conformal narrows conditional coverage from [62.9%,98.6%] to [90%,96.2%].",
        C_RED,
        RGBColor(0xFC, 0xE8, 0xE8),
    ),
]

for i, (rq_id, title, answer, col, bg) in enumerate(rq_answers):
    row, c = divmod(i, 2)
    lx = 0.3 + c * 6.55
    ty = 1.2 + row * 2.85
    add_rect(s, lx, ty, 6.3, 2.65, fill_rgb=bg, line_rgb=col, line_width=Pt(2))
    add_rect(s, lx, ty, 6.3, 0.42, fill_rgb=col)
    add_textbox(
        s,
        lx + 0.1,
        ty + 0.06,
        1.0,
        0.3,
        rq_id,
        font_size=Pt(14),
        bold=True,
        color=C_WHITE,
    )
    add_textbox(
        s,
        lx + 1.2,
        ty + 0.06,
        5.0,
        0.3,
        title,
        font_size=Pt(13),
        bold=True,
        color=C_WHITE,
    )
    add_textbox(
        s,
        lx + 0.15,
        ty + 0.5,
        6.0,
        2.05,
        answer,
        font_size=Pt(11),
        color=C_BLACK,
        word_wrap=True,
    )

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Limitations
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Limitations", "Honest assessment of what this thesis does NOT show", 28)

footer(s, slide_num=28)


limitations = [
    (
        "10% Data Subset",
        "Only 1,000 of 10,000 available scenarios used. "
        "Natterer et al. achieved R²=0.91 with full dataset. "
        "Gap from R²=0.60 likely due to 10× less training data.",
        C_ORANGE,
    ),
    (
        "Single Road Network",
        "All experiments on Paris only. "
        "Generalisation to other cities, topologies, or policy types not tested. "
        "Feature correlation patterns may differ elsewhere.",
        C_RED,
    ),
    (
        "Identical Architecture Ensembles",
        "Ensemble experiments use only PointNetTransfGAT variants. "
        "True architectural diversity (GAT + GraphSAGE + GCN) not tested — "
        "may unlock further UQ benefits.",
        C_MID_BLUE,
    ),
    (
        "Epistemic Only",
        "MC Dropout captures only epistemic (model) uncertainty. "
        "Aleatoric (data/noise) uncertainty not separated. "
        "T9 heteroscedastic attempt failed (R²=0.02 — variance inflation).",
        RGBColor(0x70, 0x30, 0xA0),
    ),
    (
        "No Non-GNN Baselines",
        "No comparison to random forests, GPs, or MLP ensembles. "
        "Unclear if ρ=0.4820 is specific to GNN or achievable by simpler models.",
        C_GRAY,
    ),
]
col = 0
ty_l = [1.2, 1.2, 3.15, 3.15, 5.1]
lx_l = [0.3, 6.65, 0.3, 6.65, 0.3]
for i, (title, text, col_c) in enumerate(limitations):
    lx = lx_l[i]
    ty = ty_l[i]
    w = 12.73 if i == 4 else 6.05
    h = 1.75
    add_rect(s, lx, ty, w, h, fill_rgb=C_LIGHT_BLUE, line_rgb=col_c, line_width=Pt(1.8))
    add_textbox(
        s,
        lx + 0.12,
        ty + 0.08,
        w - 0.24,
        0.35,
        title,
        font_size=Pt(13),
        bold=True,
        color=col_c,
    )
    add_textbox(
        s,
        lx + 0.12,
        ty + 0.45,
        w - 0.24,
        1.15,
        text,
        font_size=Pt(11),
        color=C_BLACK,
        word_wrap=True,
    )

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — Future Work
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Future Work", "Six directions for extending this research", 29)

footer(s, slide_num=29)


future = [
    (
        "Full 10k Scenarios",
        "Train on complete 10,000-scenario dataset.\nAssess if R² and ρ scale with more data.",
        C_MID_BLUE,
    ),
    (
        "Architecturally Diverse Ensembles",
        "Combine GAT + GraphSAGE + GCN models.\nTrue diversity may improve UQ beyond ρ=0.4333.",
        C_GREEN,
    ),
    (
        "Conformal Risk Control",
        "Extend to arbitrary loss functions.\nTighter adaptive intervals via conformal risk control.",
        C_ORANGE,
    ),
    (
        "Epistemic-Aleatoric Decomposition",
        "Evidential deep learning approach.\nSeparate model uncertainty from data noise.",
        RGBColor(0x70, 0x30, 0xA0),
    ),
    (
        "Multi-City Validation",
        "Test on London, Berlin, Munich networks.\nCheck generalisability of UQ findings.",
        C_RED,
    ),
    (
        "Active Learning",
        "Use high-σ predictions to select\nnew MATSim scenarios for training.",
        RGBColor(0x00, 0x70, 0x70),
    ),
]
positions_fw = [
    (0.3, 1.2),
    (4.65, 1.2),
    (9.0, 1.2),
    (0.3, 3.9),
    (4.65, 3.9),
    (9.0, 3.9),
]
for (lx, ty), (title, text, col_c) in zip(positions_fw, future):
    add_rect(
        s, lx, ty, 4.1, 2.5, fill_rgb=C_LIGHT_BLUE, line_rgb=col_c, line_width=Pt(2)
    )
    add_rect(s, lx, ty, 4.1, 0.42, fill_rgb=col_c)
    add_textbox(
        s,
        lx + 0.1,
        ty + 0.06,
        3.9,
        0.3,
        title,
        font_size=Pt(12),
        bold=True,
        color=C_WHITE,
    )
    add_textbox(
        s,
        lx + 0.12,
        ty + 0.52,
        3.85,
        1.82,
        text,
        font_size=Pt(11),
        color=C_BLACK,
        word_wrap=True,
    )

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=C_DARK_BLUE)
add_rect(s, 0, 1.4, 13.33, 4.3, fill_rgb=RGBColor(0x0E, 0x25, 0x40))

add_textbox(
    s,
    0.5,
    0.2,
    12.3,
    0.9,
    "Conclusion",
    font_size=Pt(34),
    bold=True,
    color=C_WHITE,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    s,
    0.5,
    1.0,
    12.3,
    0.38,
    "ML Surrogates for Agent-Based Transport Models",
    font_size=Pt(16),
    italic=True,
    color=C_LIGHT_BLUE,
    align=PP_ALIGN.CENTER,
)

takeaways = [
    (
        "MC Dropout works.",
        "ρ=0.4820 on 3.16M predictions. Useful ranking signal with AUROC=0.76. "
        "Cost to enable: ΔR²=−0.010 (negligible).",
    ),
    (
        "Conformal prediction delivers guarantees.",
        "90.02% and 95.01% coverage achieved by construction. "
        "Adaptive variant reduces conditional disparity from [62.9%,98.6%] to [90%,96.2%].",
    ),
    (
        "Use the 3-layer hierarchy in practice.",
        "MC Dropout for ranking → temperature scaling for calibration → "
        "conformal prediction for formal guarantees. All post-hoc, no retraining.",
    ),
]
ty = 1.52
for i, (bold_part, rest) in enumerate(takeaways):
    add_rect(s, 0.4, ty, 0.5, 0.9, fill_rgb=C_MID_BLUE)
    add_textbox(
        s,
        0.45,
        ty + 0.22,
        0.4,
        0.45,
        str(i + 1),
        font_size=Pt(18),
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_rect(s, 0.95, ty, 11.95, 0.9, fill_rgb=RGBColor(0x14, 0x2D, 0x4A))
    add_textbox(
        s,
        1.1,
        ty + 0.07,
        11.65,
        0.32,
        bold_part,
        font_size=Pt(13),
        bold=True,
        color=C_LIGHT_BLUE,
    )
    add_textbox(
        s,
        1.1,
        ty + 0.4,
        11.65,
        0.45,
        rest,
        font_size=Pt(11),
        color=RGBColor(0xCC, 0xDD, 0xEE),
        word_wrap=True,
    )
    ty += 1.02

add_textbox(
    s,
    0.5,
    5.75,
    12.3,
    0.45,
    "Deploying ML surrogates responsibly in urban planning requires "
    "not just accurate predictions, but principled uncertainty quantification.",
    font_size=Pt(13),
    italic=True,
    color=RGBColor(0xAA, 0xCC, 0xEE),
    align=PP_ALIGN.CENTER,
    word_wrap=True,
)

add_textbox(
    s,
    0.5,
    6.35,
    12.3,
    0.38,
    "Thank you  |  Questions welcome",
    font_size=Pt(17),
    bold=True,
    color=C_WHITE,
    align=PP_ALIGN.CENTER,
)

# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
